import matplotlib.pyplot as plt
import ftplib
import patoolib
import os
import glob
from zipfile import ZipFile
import pandas as pd
import pyodbc
import datetime
import jdatetime
import numpy as np
from pylab import figure, clf, plot, bar, stem, xlabel, ylabel, xlim, ylim, title, grid, axes, show, legend
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import cv2





CS_Traffic_NAK_Alborz=[0.13678734,
0.1345394,
0.19775998,
0.2338002,
0.23635279,
0.23409754,
0.22508111,
0.26681313,
0.26720917,
0.27121695,
0.2714414,
0.25152279,
0.27010468,
0.300718,
0.26771107,
0.26104586,
0.24653473,
0.25184674,
0.26355545,
0.25487226,
0.19243521,
0.2533507,
0.25678935,
0.25087218,
0.24575591,
0.22962641,
0.23657359,
0.23549804,
0.26712722,
0.26592667,
0.26248125,
0.27497933,
0.26691224,
0.26590984
]

CC_NAK_Alborz=[
83.09,
83,
82.92,
83.86,
81.94,
83.21,
82.43,
83.2,
82.26,
82.68,
83.32,
82.84,
83.28,
82.91,
81.74,
82.84,
82.6,
81.66,
81.06,
81.64,
81.98,
82.36,
83.24,
83.7,
84.15,
84.43,
84.37,
85.27,
85.58,
85.29,
85.04,
84.48,
84.32,
84.65]

Week_Vec=[
'W01',
'W02',
'W03',
'W04',
'W05',
'W06',
'W07',
'W08',
'W09',
'W10',
'W11',
'W12',
'W13',
'W14',
'W15',
'W16',
'W17',
'W18',
'W19',
'W20',
'W21',
'W22',
'W23',
'W24',
'W25',
'W26',
'W27',
'W28',
'W29',
'W30',
'W31',
'W32',
'W33',
'W34']







#prs = Presentation()
#blank_slide_layout = prs.slide_layouts[6]
#slide = prs.slides.add_slide(blank_slide_layout)

### creating the object of the Presentation class  
##myPPT = Presentation()  
  
### creating the slide layout  
##firstLayout = myPPT.slide_layouts[0]  
  
### creating the slide object to add in PPT file  
##mySlide = myPPT.slides.add_slide(firstLayout)  
  
## adding the title in the slide  
#myTitle = slide.shapes.title  
## adding the subtitle in the slide  
##mySubtitle = slide.shapes.placeholders[1]  
  
## inserting text in the title and subtitle  
#myTitle.text = "My First Presentation"  
##mySubtitle.text = "Using the python-pptx library"  
  
## saving the PPT file  
#myPPT.save('myPPT.pptx')  

#r=0

#def cm_to_inch(value):
#    return value/2.54

#x = np.arange(len(CS_Traffic_NAK_Alborz))
#fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
#ax2 = ax1.twinx()
#ax1.yaxis.tick_right()
#ax2.yaxis.tick_left()
#ax1.bar(Week_Vec,CS_Traffic_NAK_Alborz,color = "bisque")
#ax2.plot(Week_Vec,CC_NAK_Alborz,color = "darkorange")
#ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
#font1 = {'family':'serif','color':'black','size':8}
#plt.title("CC-NAK-Alborz", fontdict = font1)
#ax1.legend(['CC(%)    Tatal Traffic (MErlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
#ax1.yaxis.set_tick_params(labelsize=7)
#ax2.yaxis.set_tick_params(labelsize=7)
#plt.savefig('CC_NAK-Alborz.png')

#image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png")
#y=30
#x=5
#h=1000
#w=550
#CC_NAK_Alborz_Cropped = image[x:w, y:h]
#cv2.imwrite("CC_NAK-Alborz.png", CC_NAK_Alborz_Cropped)




#pic_left_1  = int(prs.slide_width *0)
#pic_top_1   = int(prs.slide_width *0.02)
#pic_width_1 = int(prs.slide_width *0.5)

#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png", pic_left_1, pic_top_1, pic_width_1)

#pic_left_1  = int(prs.slide_width *0.5)
#pic_top_1   = int(prs.slide_width *0.02)
#pic_width_1 = int(prs.slide_width *0.5)

#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png", pic_left_1, pic_top_1, pic_width_1)


#pic_left_1  = int(prs.slide_width *0)
#pic_top_1   = int(prs.slide_width *0.29)
#pic_width_1 = int(prs.slide_width *0.5)

#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png", pic_left_1, pic_top_1, pic_width_1)

#pic_left_1  = int(prs.slide_width *0.5)
#pic_top_1   = int(prs.slide_width *0.29)
#pic_width_1 = int(prs.slide_width *0.5)

#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png", pic_left_1, pic_top_1, pic_width_1)




#x = np.arange(len(CS_Traffic_NAK_Alborz))
##fig, ax1 = plt.subplots(figsize=(cm_to_inch(8.5),cm_to_inch(5.5)))
#fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
#ax2 = ax1.twinx()
#ax1.yaxis.tick_right()
#ax2.yaxis.tick_left()
#ax1.bar(Week_Vec,CS_Traffic_NAK_Alborz,color = "bisque")
#ax2.plot(Week_Vec,CC_NAK_Alborz,color = "darkorange")
#ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
#font1 = {'family':'serif','color':'black','size':6}
#plt.title("CC-NAK-Alborz", fontdict = font1)
#ax1.legend(['CC(%)    Tatal Traffic (MErlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
#ax1.yaxis.set_tick_params(labelsize=6)
#ax2.yaxis.set_tick_params(labelsize=6)
#plt.savefig('CC_NAK-Alborz.png')

#image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png")
#y=20
#x=5
#h=1000
#w=550
#CC_NAK_Alborz_Cropped = image[x:w, y:h]
#cv2.imwrite("CC_NAK-Alborz.png", CC_NAK_Alborz_Cropped)




#pic_left_1  = int(prs.slide_width *0)
#pic_top_1   = int(prs.slide_width *0.56)
#pic_width_1 = int(prs.slide_width *0.33)

#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png", pic_left_1, pic_top_1, pic_width_1)

#pic_left_1  = int(prs.slide_width *0.33)
#pic_top_1   = int(prs.slide_width *0.56)
#pic_width_1 = int(prs.slide_width *0.33)

#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png", pic_left_1, pic_top_1, pic_width_1)

#pic_left_1  = int(prs.slide_width *0.67)
#pic_top_1   = int(prs.slide_width *0.56)
#pic_width_1 = int(prs.slide_width *0.33)

#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_NAK-Alborz.png", pic_left_1, pic_top_1, pic_width_1)


#prs.save('test.pptx')



#plt.show()


# ****************************************************************************************************
# ((((((((((((((((((((((((((((((((( Connection to PERFORMANCEDB01 ))))))))))))))))))))))))))))))))))))
# ****************************************************************************************************
conn = pyodbc.connect('Driver={SQL Server};'
                      'Server=PERFORMANCEDB01;'
                      'Database=Performance_NAK;'
                      'Trusted_Connection=yes;')
conn_performanceDB = conn.cursor()



# /\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\/\
# ((((((((((((((((((((( Functions ))))))))))))))))))))))))
# \/\/\/\/\/\/\/\/\/\/\/\\/\/\/\/\/\/\/\/\/\/\/\\/\/\/\/\/

# This function is used to covert centimeter to inch
def cm_to_inch(value):
    return value/2.54


# This function is used to reduce sample rate with downsample Rate=Rate
def downsample(Vector,Rate):
    downsample_vec=[]
    first_Index=0
    downsample_vec.append(Vector[0])
    for k in range(len(Vector)):
        if (k-first_Index==Rate):
            downsample_vec.append(Vector[k])
            first_Index=k
    return downsample_vec

















# ****************************************************************************************************
# (((((((((((((((((((((((((((( Total Traffic and Payload per Contractor )))))))))))))))))))))))))))))))
# ****************************************************************************************************


conn_performanceDB.execute("select Wk,Contractor, sum([Total Voice Traffic (Erlang)]) as 'Total Voice Traffic (Erlang)' from ("+
                           "select Wk, Contractor,PIndex, avg([Total Voice Traffic (Erlang)]) as 'Total Voice Traffic (Erlang)' from "+
                            "Province_KPI_Score_Band_CS_Daily group by Wk, Contractor, PIndex ) tble group by Wk, Contractor  order by Wk")
CS_Country_Table=conn_performanceDB.fetchall()


CS_Traffic_NAK_Alborz=[]
CS_Traffic_NAK_Tehran=[]
CS_Traffic_NAK_North=[]
CS_Traffic_NAK_Nokia=[]
CS_Traffic_NAK_Huawei=[]
CS_Traffic_Farafan=[]
CS_Traffic_BR_TEL=[]
CS_Traffic_Huawei=[]


# ----------------------- CS ----------------------------------

for i in range(len(CS_Country_Table)):


    Row_Data=str(CS_Country_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Contractor=Row_Data[1]
    CS_Traffic=Row_Data[2]
    Week=Week[2:9]
    Contractor=Contractor[1:len(Contractor)-1]
    CS_Traffic_Val=round(float(CS_Traffic[0:len(CS_Traffic)-1])/1e6,3)

    #if Week=='1401-33':
    #    break



    if (Contractor=='NAK-Alborz'):
        CS_Traffic_NAK_Alborz.append(CS_Traffic_Val)
    if (Contractor=='NAK-Tehran'):
        CS_Traffic_NAK_Tehran.append(CS_Traffic_Val)
    if (Contractor=='NAK-North'):
        CS_Traffic_NAK_North.append(CS_Traffic_Val)
    if (Contractor=='NAK-Nokia'):
        CS_Traffic_NAK_Nokia.append(CS_Traffic_Val)
    if (Contractor=='NAK-Huawei'):
        CS_Traffic_NAK_Huawei.append(CS_Traffic_Val)
    if (Contractor=='Farafan'):
        CS_Traffic_Farafan.append(CS_Traffic_Val)
    if (Contractor=='BR-TEL'):
        CS_Traffic_BR_TEL.append(CS_Traffic_Val)
    if (Contractor=='Huawei'):
        CS_Traffic_Huawei.append(CS_Traffic_Val)


# Sort Data Based on Last Values
Last_CS_Traffic_Value=[CS_Traffic_NAK_Alborz[len(CS_Traffic_NAK_Alborz)-1], CS_Traffic_NAK_Tehran[len(CS_Traffic_NAK_Tehran)-1], CS_Traffic_NAK_North[len(CS_Traffic_NAK_North)-1],  CS_Traffic_NAK_Nokia[len(CS_Traffic_NAK_Nokia)-1], CS_Traffic_NAK_Huawei[len(CS_Traffic_NAK_Huawei)-1], CS_Traffic_Farafan[len(CS_Traffic_Farafan)-1], CS_Traffic_BR_TEL[len(CS_Traffic_BR_TEL)-1],  CS_Traffic_Huawei[len(CS_Traffic_Huawei)-1]]
Index_of_Sort=np.argsort(Last_CS_Traffic_Value)

Data_Sorted_Array=[]
x_Labels=[];
for k in range(len(Index_of_Sort)):
    if Index_of_Sort[k]==0:
        Data_Sorted_Array.append(CS_Traffic_NAK_Alborz)
        x_Labels.append('NAK-Alborz')
    if Index_of_Sort[k]==1:
        Data_Sorted_Array.append(CS_Traffic_NAK_Tehran)
        x_Labels.append('NAK-Tehran')
    if Index_of_Sort[k]==2:
        Data_Sorted_Array.append(CS_Traffic_NAK_North)
        x_Labels.append('NAK-North')
    if Index_of_Sort[k]==3:
        Data_Sorted_Array.append(CS_Traffic_NAK_Nokia)
        x_Labels.append('NAK-Nokia')
    if Index_of_Sort[k]==4:
        Data_Sorted_Array.append(CS_Traffic_NAK_Huawei)
        x_Labels.append('NAK-Huawei')
    if Index_of_Sort[k]==5:
        Data_Sorted_Array.append(CS_Traffic_Farafan)
        x_Labels.append('Farafan')
    if Index_of_Sort[k]==6:
        Data_Sorted_Array.append(CS_Traffic_BR_TEL)
        x_Labels.append('BR_TEL')
    if Index_of_Sort[k]==7:
        Data_Sorted_Array.append(CS_Traffic_Huawei)
        x_Labels.append('Huawei')



Last_CS_Traffic_Value_NAK=[CS_Traffic_NAK_Alborz[len(CS_Traffic_NAK_Alborz)-1]+ CS_Traffic_NAK_Tehran[len(CS_Traffic_NAK_Tehran)-1]+ CS_Traffic_NAK_North[len(CS_Traffic_NAK_North)-1]+  CS_Traffic_NAK_Nokia[len(CS_Traffic_NAK_Nokia)-1]+ CS_Traffic_NAK_Huawei[len(CS_Traffic_NAK_Huawei)-1], CS_Traffic_Farafan[len(CS_Traffic_Farafan)-1], CS_Traffic_BR_TEL[len(CS_Traffic_BR_TEL)-1],  CS_Traffic_Huawei[len(CS_Traffic_Huawei)-1]]
Index_of_Sort_NAK=np.argsort(Last_CS_Traffic_Value_NAK)
Data_Sorted_Array_Nak=[]
x_Labels_Nak=[];
for k in range(len(Index_of_Sort_NAK)):
    if Index_of_Sort_NAK[k]==0:
        A1=np.add(CS_Traffic_NAK_Alborz,CS_Traffic_NAK_Tehran)
        A2=np.add(A1,CS_Traffic_NAK_North)
        A3=np.add(A2,CS_Traffic_NAK_Nokia)
        Data_Sorted_Array_Nak.append(np.add(A3,CS_Traffic_NAK_Huawei))
        x_Labels_Nak.append('NAK')
    if Index_of_Sort_NAK[k]==1:
        Data_Sorted_Array_Nak.append(CS_Traffic_Farafan)
        x_Labels_Nak.append('Farafan')
    if Index_of_Sort_NAK[k]==2:
        Data_Sorted_Array_Nak.append(CS_Traffic_BR_TEL)
        x_Labels_Nak.append('BR_TEL')
    if Index_of_Sort_NAK[k]==3:
        Data_Sorted_Array_Nak.append(CS_Traffic_Huawei)
        x_Labels_Nak.append('Huawei')


data=np.array(Data_Sorted_Array)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(12)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_CS_Traffic_Value):
    plt.text( i + dx[31],Last_CS_Traffic_Value[Index_of_Sort[i]] , str(Last_CS_Traffic_Value[Index_of_Sort[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels)
font1 = {'family':'serif','color':'black','size':17}
plt.title("Total Traffic (MErlang)", fontdict = font1)
grid(True)
plt.savefig('CS_Traffic_Bar.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CS_Traffic_Bar.png")
y=80
x=20
h=1000
w=520
CS_Traffic_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("CS_Traffic_Bar.png", CS_Traffic_Bar_Cropped)




Data_Sorted_Pie=[Data_Sorted_Array[0][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[1][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[2][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[3][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[4][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[5][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[6][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[7][len(Data_Sorted_Array[0])-1]]
plt.figure(figsize=(cm_to_inch(12),cm_to_inch(10)))
#font1 = {'family':'serif','color':'black','size':15}
#plt.title("Total Traffic (%)", fontdict = font1)
plt.pie(Data_Sorted_Pie,labels =x_Labels,autopct='%1.1f%%')
plt.savefig('CS_Traffic_Bar_Percentage_1.png')


image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CS_Traffic_Bar_Percentage_1.png")
y=50
x=20
h=450
w=350
CS_Traffic_Bar_Percentage_1_Cropped = image[x:w, y:h]
cv2.imwrite("CS_Traffic_Bar_Percentage_1.png", CS_Traffic_Bar_Percentage_1_Cropped)


Data_Sorted_NAK_Pie=[Data_Sorted_Array_Nak[0][len(Data_Sorted_Array_Nak[0])-1], Data_Sorted_Array_Nak[1][len(Data_Sorted_Array_Nak[0])-1], Data_Sorted_Array_Nak[2][len(Data_Sorted_Array_Nak[0])-1], Data_Sorted_Array_Nak[3][len(Data_Sorted_Array_Nak[0])-1]]
plt.figure(figsize=(cm_to_inch(12),cm_to_inch(10)))
#font1 = {'family':'serif','color':'black','size':15}
#plt.title("Total Traffic (%)", fontdict = font1)
plt.pie(Data_Sorted_NAK_Pie,labels =x_Labels_Nak,autopct='%1.1f%%')
plt.savefig('CS_Traffic_Bar_Percentage_2.png')
#plt.show() 

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CS_Traffic_Bar_Percentage_2.png")
y=50
x=20
h=450
w=350
CS_Traffic_Bar_Percentage_2_Cropped = image[x:w, y:h]
cv2.imwrite("CS_Traffic_Bar_Percentage_2.png", CS_Traffic_Bar_Percentage_2_Cropped)




prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

#pic_width_NAK = int(prs.slide_width *0.1)
#pic_left_MCI = int(prs.slide_width *0.9)
#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\\NAK.png", 0, 0, pic_width_NAK)
#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\\MCI.png", pic_left_MCI, 0, pic_width_NAK)

pic_left_1  = int(prs.slide_width *0.045)
pic_top_1   = 0
pic_width_1 = int(prs.slide_width *0.9)

pic_left_2  = int(prs.slide_width *0.1)
pic_top_2   = int(prs.slide_width *0.44)
pic_width_2 = int(prs.slide_width *0.37)

pic_left_3  = int(prs.slide_width *0.55)
pic_top_3   = int(prs.slide_width *0.44)
pic_width_3 = int(prs.slide_width *0.37)

pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CS_Traffic_Bar.png", pic_left_1, pic_top_1, pic_width_1)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CS_Traffic_Bar_Percentage_1.png", pic_left_2, pic_top_2, pic_width_2)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CS_Traffic_Bar_Percentage_2.png", pic_left_3, pic_top_3, pic_width_3)
prs.save('test.pptx')







conn_performanceDB.execute("select Wk,Contractor, sum([Total Payload (GB)]) as 'Total Payload (GB)' from ("+
"select Wk, Contractor,[Province Index], avg([Total Payload (GB)]) as 'Total Payload (GB)' from "+
"Province_KPI_Score_Band_PS_Daily group by Wk, Contractor, [Province Index] ) tble group by Wk, Contractor  order by Wk")
PS_Country_Table=conn_performanceDB.fetchall()


PS_Traffic_NAK_Alborz=[]
PS_Traffic_NAK_Tehran=[]
PS_Traffic_NAK_North=[]
PS_Traffic_NAK_Nokia=[]
PS_Traffic_NAK_Huawei=[]
PS_Traffic_Farafan=[]
PS_Traffic_BR_TEL=[]
PS_Traffic_Huawei=[]


# ----------------------- PS ----------------------------------

for i in range(len(PS_Country_Table)):


    Row_Data=str(PS_Country_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Contractor=Row_Data[1]
    PS_Traffic=Row_Data[2]
    Week=Week[2:9]
    Contractor=Contractor[1:len(Contractor)-1]
    PS_Traffic_Val=round(float(PS_Traffic[0:len(PS_Traffic)-1])/1e6,3)

    #if Week=='1401-33':
    #    break



    if (Contractor=='NAK-Alborz'):
        PS_Traffic_NAK_Alborz.append(PS_Traffic_Val)
    if (Contractor=='NAK-Tehran'):
        PS_Traffic_NAK_Tehran.append(PS_Traffic_Val)
    if (Contractor=='NAK-North'):
        PS_Traffic_NAK_North.append(PS_Traffic_Val)
    if (Contractor=='NAK-Nokia'):
        PS_Traffic_NAK_Nokia.append(PS_Traffic_Val)
    if (Contractor=='NAK-Huawei'):
        PS_Traffic_NAK_Huawei.append(PS_Traffic_Val)
    if (Contractor=='Farafan'):
        PS_Traffic_Farafan.append(PS_Traffic_Val)
    if (Contractor=='BR-TEL'):
        PS_Traffic_BR_TEL.append(PS_Traffic_Val)
    if (Contractor=='Huawei'):
        PS_Traffic_Huawei.append(PS_Traffic_Val)


# Sort Data Based on Last Values
Last_PS_Traffic_Value=[PS_Traffic_NAK_Alborz[len(PS_Traffic_NAK_Alborz)-1], PS_Traffic_NAK_Tehran[len(PS_Traffic_NAK_Tehran)-1], PS_Traffic_NAK_North[len(PS_Traffic_NAK_North)-1],  PS_Traffic_NAK_Nokia[len(PS_Traffic_NAK_Nokia)-1], PS_Traffic_NAK_Huawei[len(PS_Traffic_NAK_Huawei)-1], PS_Traffic_Farafan[len(PS_Traffic_Farafan)-1], PS_Traffic_BR_TEL[len(PS_Traffic_BR_TEL)-1],  PS_Traffic_Huawei[len(PS_Traffic_Huawei)-1]]
Index_of_Sort=np.argsort(Last_PS_Traffic_Value)

Data_Sorted_Array=[]
x_Labels=[];
for k in range(len(Index_of_Sort)):
    if Index_of_Sort[k]==0:
        Data_Sorted_Array.append(PS_Traffic_NAK_Alborz)
        x_Labels.append('NAK-Alborz')
    if Index_of_Sort[k]==1:
        Data_Sorted_Array.append(PS_Traffic_NAK_Tehran)
        x_Labels.append('NAK-Tehran')
    if Index_of_Sort[k]==2:
        Data_Sorted_Array.append(PS_Traffic_NAK_North)
        x_Labels.append('NAK-North')
    if Index_of_Sort[k]==3:
        Data_Sorted_Array.append(PS_Traffic_NAK_Nokia)
        x_Labels.append('NAK-Nokia')
    if Index_of_Sort[k]==4:
        Data_Sorted_Array.append(PS_Traffic_NAK_Huawei)
        x_Labels.append('NAK-Huawei')
    if Index_of_Sort[k]==5:
        Data_Sorted_Array.append(PS_Traffic_Farafan)
        x_Labels.append('Farafan')
    if Index_of_Sort[k]==6:
        Data_Sorted_Array.append(PS_Traffic_BR_TEL)
        x_Labels.append('BR_TEL')
    if Index_of_Sort[k]==7:
        Data_Sorted_Array.append(PS_Traffic_Huawei)
        x_Labels.append('Huawei')



Last_PS_Traffic_Value_NAK=[PS_Traffic_NAK_Alborz[len(PS_Traffic_NAK_Alborz)-1]+ PS_Traffic_NAK_Tehran[len(PS_Traffic_NAK_Tehran)-1]+ PS_Traffic_NAK_North[len(PS_Traffic_NAK_North)-1]+  PS_Traffic_NAK_Nokia[len(PS_Traffic_NAK_Nokia)-1]+ PS_Traffic_NAK_Huawei[len(PS_Traffic_NAK_Huawei)-1], PS_Traffic_Farafan[len(PS_Traffic_Farafan)-1], PS_Traffic_BR_TEL[len(PS_Traffic_BR_TEL)-1],  PS_Traffic_Huawei[len(PS_Traffic_Huawei)-1]]
Index_of_Sort_NAK=np.argsort(Last_PS_Traffic_Value_NAK)
Data_Sorted_Array_Nak=[]
x_Labels_Nak=[];
for k in range(len(Index_of_Sort_NAK)):
    if Index_of_Sort_NAK[k]==0:
        A1=np.add(PS_Traffic_NAK_Alborz,PS_Traffic_NAK_Tehran)
        A2=np.add(A1,PS_Traffic_NAK_North)
        A3=np.add(A2,PS_Traffic_NAK_Nokia)
        Data_Sorted_Array_Nak.append(np.add(A3,PS_Traffic_NAK_Huawei))
        x_Labels_Nak.append('NAK')
    if Index_of_Sort_NAK[k]==1:
        Data_Sorted_Array_Nak.append(PS_Traffic_Farafan)
        x_Labels_Nak.append('Farafan')
    if Index_of_Sort_NAK[k]==2:
        Data_Sorted_Array_Nak.append(PS_Traffic_BR_TEL)
        x_Labels_Nak.append('BR_TEL')
    if Index_of_Sort_NAK[k]==3:
        Data_Sorted_Array_Nak.append(PS_Traffic_Huawei)
        x_Labels_Nak.append('Huawei')


data=np.array(Data_Sorted_Array)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(12)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_PS_Traffic_Value):
    plt.text( i + dx[31],Last_PS_Traffic_Value[Index_of_Sort[i]] , str(Last_PS_Traffic_Value[Index_of_Sort[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels)
font1 = {'family':'serif','color':'black','size':17}
plt.title("Total Payload (PB)", fontdict = font1)
grid(True)
plt.savefig('PS_Traffic_Bar.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\PS_Traffic_Bar.png")
y=80
x=20
h=1000
w=520
PS_Traffic_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("PS_Traffic_Bar.png", PS_Traffic_Bar_Cropped)




Data_Sorted_Pie=[Data_Sorted_Array[0][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[1][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[2][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[3][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[4][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[5][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[6][len(Data_Sorted_Array[0])-1], Data_Sorted_Array[7][len(Data_Sorted_Array[0])-1]]
plt.figure(figsize=(cm_to_inch(12),cm_to_inch(10)))
#font1 = {'family':'serif','color':'black','size':15}
#plt.title("Total Traffic (%)", fontdict = font1)
plt.pie(Data_Sorted_Pie,labels =x_Labels,autopct='%1.1f%%')
plt.savefig('PS_Traffic_Bar_Percentage_1.png')


image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\PS_Traffic_Bar_Percentage_1.png")
y=50
x=20
h=450
w=350
PS_Traffic_Bar_Percentage_1_Cropped = image[x:w, y:h]
cv2.imwrite("PS_Traffic_Bar_Percentage_1.png", PS_Traffic_Bar_Percentage_1_Cropped)


Data_Sorted_NAK_Pie=[Data_Sorted_Array_Nak[0][len(Data_Sorted_Array_Nak[0])-1], Data_Sorted_Array_Nak[1][len(Data_Sorted_Array_Nak[0])-1], Data_Sorted_Array_Nak[2][len(Data_Sorted_Array_Nak[0])-1], Data_Sorted_Array_Nak[3][len(Data_Sorted_Array_Nak[0])-1]]
plt.figure(figsize=(cm_to_inch(12),cm_to_inch(10)))
#font1 = {'family':'serif','color':'black','size':15}
#plt.title("Total Traffic (%)", fontdict = font1)
plt.pie(Data_Sorted_NAK_Pie,labels =x_Labels_Nak,autopct='%1.1f%%')
plt.savefig('PS_Traffic_Bar_Percentage_2.png')
#plt.show() 

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\PS_Traffic_Bar_Percentage_2.png")
y=50
x=20
h=450
w=350
PS_Traffic_Bar_Percentage_2_Cropped = image[x:w, y:h]
cv2.imwrite("PS_Traffic_Bar_Percentage_2.png", PS_Traffic_Bar_Percentage_2_Cropped)




slide = prs.slides.add_slide(blank_slide_layout)

#pic_width_NAK = int(prs.slide_width *0.1)
#pic_left_MCI = int(prs.slide_width *0.9)
#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\\NAK.png", 0, 0, pic_width_NAK)
#pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\\MCI.png", pic_left_MCI, 0, pic_width_NAK)

pic_left_1  = int(prs.slide_width *0.045)
pic_top_1   = 0
pic_width_1 = int(prs.slide_width *0.9)

pic_left_2  = int(prs.slide_width *0.1)
pic_top_2   = int(prs.slide_width *0.44)
pic_width_2 = int(prs.slide_width *0.37)

pic_left_3  = int(prs.slide_width *0.55)
pic_top_3   = int(prs.slide_width *0.44)
pic_width_3 = int(prs.slide_width *0.37)

pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\PS_Traffic_Bar.png", pic_left_1, pic_top_1, pic_width_1)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\PS_Traffic_Bar_Percentage_1.png", pic_left_2, pic_top_2, pic_width_2)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\PS_Traffic_Bar_Percentage_2.png", pic_left_3, pic_top_3, pic_width_3)
prs.save('test.pptx')





# ****************************************************************************************************
# (((((((((((((((((((((( Total Traffic and Payload per Contractor and Band ))))))))))))))))))))))))))))
# ****************************************************************************************************

for t in range(16):
        if (t==0 or t==8):
            Contractor="NAK-Alborz"
        if (t==1 or t==9):
            Contractor="NAK-North"
        if (t==2 or t==10):
            Contractor="NAK-Tehran"
        if (t==3 or t==11):
            Contractor="NAK-Huawei"
        if (t==4 or t==12):
            Contractor="NAK-Nokia"
        if (t==5 or t==13):
            Contractor="BR-TEL"
        if (t==6 or t==14):
            Contractor="Farafan"
        if (t==7 or t==15):
            Contractor="Huawei"



        Traffic_2G=[]
        Traffic_U900F1=[]
        Traffic_U900F2=[]
        Traffic_U2100F1=[]
        Traffic_U2100F2=[]
        Traffic_U2100F3=[]
        Traffic_L1800F1=[]
        Traffic_L1800F2=[]
        Traffic_L2100F1=[]
        Traffic_L2100F2=[]
        Traffic_L2600F1=[]
        Traffic_L2600F2=[]
        Traffic_L2300F1=[]
        Traffic_L2300F2=[]

        if t<=7:
            conn_performanceDB.execute("select Wk,Contractor, sum([2G TCH Traffic]) as '2G TCH Traffic', sum([U900F1 Traffic]) as 'U900F1 Traffic', sum([U900F2 Traffic]) as 'U900F2 Traffic', sum([U2100F1 Traffic]) as 'U2100F1 Traffic', sum([U2100F2 Traffic]) as 'U2100F2 Traffic', sum([U2100F3 Traffic]) as 'U2100F3 Traffic',"+
            " sum([L1800F1 Traffic]) as 'L1800F1 Traffic', sum([L1800F2 Traffic]) as 'L1800F2 Traffic', sum([L2100F1 Traffic]) as 'L2100F1 Traffic', sum([L2100F2 Traffic]) as 'L2100F2 Traffic', sum([L2600F1 Traffic]) as 'L2600F1 Traffic', sum([L2600F2 Traffic]) as 'L2600F2 Traffic', sum([L2300F1 Traffic]) as 'L2300F1 Traffic', sum([L2300F2 Traffic]) as 'L2300F2 Traffic'"+
            " from (select Wk, Contractor,PIndex,  avg([2G TCH Traffic]) as '2G TCH Traffic', avg([U900F1 Traffic]) as 'U900F1 Traffic', avg([U900F2 Traffic]) as 'U900F2 Traffic', avg([U2100F1 Traffic]) as 'U2100F1 Traffic', avg([U2100F2 Traffic]) as 'U2100F2 Traffic', avg([U2100F3 Traffic]) as 'U2100F3 Traffic'," +
            " avg([L1800F1 Traffic]) as 'L1800F1 Traffic', avg([L1800F2 Traffic]) as 'L1800F2 Traffic', avg([L2100F1 Traffic]) as 'L2100F1 Traffic', avg([L2100F2 Traffic]) as 'L2100F2 Traffic', avg([L2600F1 Traffic]) as 'L2600F1 Traffic', avg([L2600F2 Traffic]) as 'L2600F2 Traffic', avg([L2300F1 Traffic]) as 'L2300F1 Traffic', avg([L2300F2 Traffic]) as 'L2300F2 Traffic'"+
            " from Province_KPI_Score_Band_CS_Daily group by Wk, Contractor, PIndex ) tble where Contractor='"+Contractor+"' group by Wk, Contractor  order by Wk")
            Contractor_Table=conn_performanceDB.fetchall()
        else:
            conn_performanceDB.execute("select Wk,Contractor, sum([2G Payload]) as '2G Payload', sum([U900F1 Payload]) as 'U900F1 Payload', sum([U900F2 Payload]) as 'U900F2 Payload', sum([U2100F1 Payload]) as 'U2100F1 Payload', sum([U2100F2 Payload]) as 'U2100F2 Payload', sum([U2100F3 Payload]) as 'U2100F3 Payload',"+
            " sum([L1800F1 Payload]) as 'L1800F1 Payload', sum([L1800F2 Payload]) as 'L1800F2 Payload', sum([L2100F1 Payload]) as 'L2100F1 Payload', sum([L2100F2 Payload]) as 'L2100F2 Payload', sum([L2600F1 Payload]) as 'L2600F1 Payload', sum([L2600F2 Payload]) as 'L2600F2 Payload', sum([L2300F1 Payload]) as 'L2300F1 Payload', sum([L2300F2 Payload]) as 'L2300F2 Payload'"+
            " from (select Wk, Contractor,[Province Index],  avg([2G PS Traffic (GB)]) as '2G Payload', avg([U900F1 Payload]) as 'U900F1 Payload', avg([U900F2 Payload]) as 'U900F2 Payload', avg([U2100F1 Payload]) as 'U2100F1 Payload', avg([U2100F2 Payload]) as 'U2100F2 Payload', avg([U2100F3 Payload]) as 'U2100F3 Payload'," +
            " avg([L1800F1 Payload]) as 'L1800F1 Payload', avg([L1800F2 Payload]) as 'L1800F2 Payload', avg([L2100F1 Payload]) as 'L2100F1 Payload', avg([L2100F2 Payload]) as 'L2100F2 Payload', avg([L2600F1 Payload]) as 'L2600F1 Payload', avg([L2600F2 Payload]) as 'L2600F2 Payload', avg([L2300F1 Payload]) as 'L2300F1 Payload', avg([L2300F2 Payload]) as 'L2300F2 Payload'"+
            " from Province_KPI_Score_Band_PS_Daily group by Wk, Contractor, [Province Index] ) tble where Contractor='"+Contractor+"' group by Wk, Contractor  order by Wk")
            Contractor_Table=conn_performanceDB.fetchall()


        for i in range(len(Contractor_Table)):


            Row_Data=str(Contractor_Table[i])
            Row_Data=Row_Data.split(", ")

            Week=Row_Data[0]
            Contractor=Row_Data[1]
            Traffic_2G_Str=Row_Data[2]
            Traffic_2G_Val=round(float(Traffic_2G_Str[0:len(Traffic_2G_Str)-1])/1e3,2)
            Traffic_U900F1_Str=Row_Data[3]
            Traffic_U900F1_Val=round(float(Traffic_U900F1_Str[0:len(Traffic_U900F1_Str)-1])/1e3,2)
            Traffic_U900F2_Str=Row_Data[4]
            Traffic_U900F2_Val=round(float(Traffic_U900F2_Str[0:len(Traffic_U900F2_Str)-1])/1e3,2)
            Traffic_U2100F1_Str=Row_Data[5]
            Traffic_U2100F1_Val=round(float(Traffic_U2100F1_Str[0:len(Traffic_U2100F1_Str)-1])/1e3,2)
            Traffic_U2100F2_Str=Row_Data[6]
            Traffic_U2100F2_Val=round(float(Traffic_U2100F2_Str[0:len(Traffic_U2100F2_Str)-1])/1e3,2)
            Traffic_U2100F3_Str=Row_Data[7]
            Traffic_U2100F3_Val=round(float(Traffic_U2100F3_Str[0:len(Traffic_U2100F3_Str)-1])/1e3,2)
            Traffic_L1800F1_Str=Row_Data[8]
            Traffic_L1800F1_Val=round(float(Traffic_L1800F1_Str[0:len(Traffic_L1800F1_Str)-1])/1e3,2)
            Traffic_L1800F2_Str=Row_Data[9]
            Traffic_L1800F2_Val=round(float(Traffic_L1800F2_Str[0:len(Traffic_L1800F2_Str)-1])/1e3,2)
            Traffic_L2100F1_Str=Row_Data[10]
            Traffic_L2100F1_Val=round(float(Traffic_L2100F1_Str[0:len(Traffic_L2100F1_Str)-1])/1e3,2)
            Traffic_L2100F2_Str=Row_Data[11]
            Traffic_L2100F2_Val=round(float(Traffic_L2100F2_Str[0:len(Traffic_L2100F2_Str)-1])/1e3,2)
            Traffic_L2600F1_Str=Row_Data[12]
            Traffic_L2600F1_Val=round(float(Traffic_L2600F1_Str[0:len(Traffic_L2600F1_Str)-1])/1e3,2)
            Traffic_L2600F2_Str=Row_Data[13]
            Traffic_L2600F2_Val=round(float(Traffic_L2600F2_Str[0:len(Traffic_L2600F2_Str)-1])/1e3,2)
            Traffic_L2300F1_Str=Row_Data[14]
            Traffic_L2300F1_Val=round(float(Traffic_L2300F1_Str[0:len(Traffic_L2300F1_Str)-1])/1e3,2)
            Traffic_L2300F2_Str=Row_Data[15]
            Traffic_L2300F2_Val=round(float(Traffic_L2300F2_Str[0:len(Traffic_L2300F2_Str)-1])/1e3,2)


            Week=Week[2:9]
            Contractor=Contractor[1:len(Contractor)-1]
    

            #if Week=='1401-33':
            #    break

            Traffic_2G.append(Traffic_2G_Val)
            Traffic_U900F1.append(Traffic_U900F1_Val)
            Traffic_U900F2.append(Traffic_U900F2_Val)
            Traffic_U2100F1.append(Traffic_U2100F1_Val)
            Traffic_U2100F2.append(Traffic_U2100F2_Val)
            Traffic_U2100F3.append(Traffic_U2100F3_Val)
            Traffic_L1800F1.append(Traffic_L1800F1_Val)
            Traffic_L1800F2.append(Traffic_L1800F2_Val)
            Traffic_L2100F1.append(Traffic_L2100F1_Val)
            Traffic_L2100F2.append(Traffic_L2100F2_Val)
            Traffic_L2600F1.append(Traffic_L2600F1_Val)
            Traffic_L2600F2.append(Traffic_L2600F2_Val)
            Traffic_L2300F1.append(Traffic_L2300F1_Val)
            Traffic_L2300F2.append(Traffic_L2300F2_Val)


        # Sort Data Based on Last Values
        Last_Traffic_Value=[]
        Last_Traffic_Value=[Traffic_2G[len(Traffic_2G)-1], Traffic_U900F1[len(Traffic_2G)-1], Traffic_U900F2[len(Traffic_2G)-1],  Traffic_U2100F1[len(Traffic_2G)-1], Traffic_U2100F2[len(Traffic_2G)-1], Traffic_U2100F3[len(Traffic_2G)-1], Traffic_L1800F1[len(Traffic_2G)-1],  Traffic_L1800F2[len(Traffic_2G)-1], Traffic_L2100F1[len(Traffic_2G)-1],  Traffic_L2100F2[len(Traffic_2G)-1], Traffic_L2600F1[len(Traffic_2G)-1],  Traffic_L2600F2[len(Traffic_2G)-1], Traffic_L2300F1[len(Traffic_2G)-1],  Traffic_L2300F2[len(Traffic_2G)-1]]
        Index_of_Sort=np.argsort(Last_Traffic_Value)

        Data_Sorted_Array=[]
        x_Labels=[];
        for k in range(len(Index_of_Sort)):
            if Index_of_Sort[k]==0:
                if (np.sum(Traffic_2G)!=0):
                    Data_Sorted_Array.append(Traffic_2G)
                    x_Labels.append('2G')
            if Index_of_Sort[k]==1:
                if (np.sum(Traffic_U900F1)!=0):
                    Data_Sorted_Array.append(Traffic_U900F1)
                    x_Labels.append('U900F1')
            if Index_of_Sort[k]==2:
                if (np.sum(Traffic_U900F2)!=0):
                    Data_Sorted_Array.append(Traffic_U900F2)
                    x_Labels.append('U900F2')
            if Index_of_Sort[k]==3:
                if (np.sum(Traffic_U2100F1)!=0):
                    Data_Sorted_Array.append(Traffic_U2100F1)
                    x_Labels.append('U2100F1')
            if Index_of_Sort[k]==4:
                if (np.sum(Traffic_U2100F2)!=0):
                    Data_Sorted_Array.append(Traffic_U2100F2)
                    x_Labels.append('U2100F2')
            if Index_of_Sort[k]==5:
                if (np.sum(Traffic_U2100F3)!=0):
                    Data_Sorted_Array.append(Traffic_U2100F3)
                    x_Labels.append('U2100F3')
            if Index_of_Sort[k]==6:
                if (np.sum(Traffic_L1800F1)!=0):
                    Data_Sorted_Array.append(Traffic_L1800F1)
                    x_Labels.append('L1800F1')
            if Index_of_Sort[k]==7:
                if (np.sum(Traffic_L1800F2)!=0):
                    Data_Sorted_Array.append(Traffic_L1800F2)
                    x_Labels.append('L1800F2')
            if Index_of_Sort[k]==8:
                if (np.sum(Traffic_L2100F1)!=0):
                    Data_Sorted_Array.append(Traffic_L2100F1)
                    x_Labels.append('L2100F1')
            if Index_of_Sort[k]==9:
                if (np.sum(Traffic_L2100F2)!=0):
                    Data_Sorted_Array.append(Traffic_L2100F2)
                    x_Labels.append('L2100F2')
            if Index_of_Sort[k]==10:
                if (np.sum(Traffic_L2600F1)!=0):
                    Data_Sorted_Array.append(Traffic_L2600F1)
                    x_Labels.append('L2600F1')
            if Index_of_Sort[k]==11:
                if (np.sum(Traffic_L2600F2)!=0):
                    Data_Sorted_Array.append(Traffic_L2600F2)
                    x_Labels.append('L2600F2')
            if Index_of_Sort[k]==12:
                if (np.sum(Traffic_L2300F1)!=0):
                    Data_Sorted_Array.append(Traffic_L2300F1)
                    x_Labels.append('L2300F1')
            if Index_of_Sort[k]==13:
                if (np.sum(Traffic_L2300F2)!=0):
                    Data_Sorted_Array.append(Traffic_L2300F2)
                    x_Labels.append('L2300F2')




        data=np.array(Data_Sorted_Array)

        x = np.arange(data.shape[0])
        dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
        d = 1./(data.shape[1]+2.)

        def cm_to_inch(value):
            return value/2.54
        plt.figure(figsize=(cm_to_inch(28),cm_to_inch(12)))
        axes= plt.axes()

        if t<=7:
            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))

            for k in range(data.shape[0]):
                Last_Value=data[k,data.shape[1]-1]
                plt.text( k+ dx[31],Last_Value , str(Last_Value), color='black', size=10, fontweight='bold')
        else:
            for i in range(data.shape[1]):
                plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))

            for k in range(data.shape[0]):
                Last_Value=data[k,data.shape[1]-1]
                plt.text( k+ dx[31],Last_Value , str(Last_Value), color='black', size=10, fontweight='bold')

        #zero_count=0
        #for i , v in enumerate(Last_Traffic_Value):
        #    if Last_Traffic_Value[Index_of_Sort[i]]==0:
        #        zero_count+=1;
        #        continue
        #    plt.text( i -zero_count+ dx[31],Last_Traffic_Value[Index_of_Sort[i]] , str(Last_Traffic_Value[Index_of_Sort[i]]), color='green', size=12, fontweight='bold')

        axes.set_xticks(x, x_Labels,fontsize=7)
        #axes.set_xticks(fontsize=10)
        font1 = {'family':'serif','color':'black','size':14}
        if t<=7:
            plt.title(Contractor+" Total Traffic (KErlang)", fontdict = font1)
            grid(True)
            plt.savefig('Traffic_'+Contractor+'_Bar.png')
        else:
            plt.title(Contractor+" Total Payload (TB)", fontdict = font1)
            grid(True)
            plt.savefig('Payload_'+Contractor+'_Bar.png')



        Data_Sorted_Pie=[]
        for j in range(len(Data_Sorted_Array)):
            Data_Sorted_Pie.append(Data_Sorted_Array[j][len(Data_Sorted_Array[0])-1])

        Data_Sorted_Pie=list([Data_Sorted_Pie]/np.sum(Data_Sorted_Pie)*100)
        fig, ax = plt.subplots(figsize=(8,6))
        x = np.arange(len(Data_Sorted_Pie[0])) 
        width = 0.5
        rects1 = ax.barh(x - width/2, Data_Sorted_Pie[0], width)
        ax.bar_label(rects1,label=Data_Sorted_Pie[0],  fmt='%.2f', fontsize=15)
        ax.set_yticks(x, labels=x_Labels,fontsize=9)
        if t<=7:
            ax.set_xlabel(Contractor+' Total Traffic (%)', fontsize=14)
            plt.savefig('Traffic_'+Contractor+'_Bar_Percentage.png')
        else:
            ax.set_xlabel(Contractor+' Total Payload (%)', fontsize=14)
            plt.savefig('Payload_'+Contractor+'_Bar_Percentage.png')

        slide = prs.slides.add_slide(blank_slide_layout)

        #pic_width_NAK = int(prs.slide_width *0.1)
        #pic_left_MCI = int(prs.slide_width *0.9)
        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\\NAK.png", 0, 0, pic_width_NAK)
        #pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\\MCI.png", pic_left_MCI, 0, pic_width_NAK)

        if t<=7:
            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Traffic_"+Contractor+"_Bar.png")
            y=80
            x=20
            h=1050
            w=520
            Traffic_Contractor_Bar = image[x:w, y:h]
            cv2.imwrite("Traffic_"+Contractor+"_Bar.png", Traffic_Contractor_Bar)


            pic_left_1  = int(prs.slide_width *0.045)
            pic_top_1   = 0
            pic_width_1 = int(prs.slide_width *0.9)


            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Traffic_"+Contractor+"_Bar_Percentage.png")
            y=5
            x=40
            h=770
            w=600
            Traffic_Contractor_Bar_Percentage = image[x:w, y:h]
            cv2.imwrite("Traffic_"+Contractor+"_Bar_Percentage.png", Traffic_Contractor_Bar_Percentage)
            pic_left_3  = int(prs.slide_width *0.24)
            pic_top_3   = int(prs.slide_width *0.4)
            pic_width_3 = int(prs.slide_width *0.47)


            pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Traffic_"+Contractor+"_Bar.png", pic_left_1, pic_top_1, pic_width_1)
            pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Traffic_"+Contractor+"_Bar_Percentage.png", pic_left_3, pic_top_3, pic_width_3)

            prs.save('test.pptx')
        else:
            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Payload_"+Contractor+"_Bar.png")
            y=80
            x=20
            h=1050
            w=520
            Payload_Contractor_Bar = image[x:w, y:h]
            cv2.imwrite("Payload_"+Contractor+"_Bar.png", Payload_Contractor_Bar)


            pic_left_1  = int(prs.slide_width *0.045)
            pic_top_1   = 0
            pic_width_1 = int(prs.slide_width *0.9)


            image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Payload_"+Contractor+"_Bar_Percentage.png")
            y=5
            x=40
            h=770
            w=600
            Payload_Contractor_Bar_Percentage = image[x:w, y:h]
            cv2.imwrite("Payload_"+Contractor+"_Bar_Percentage.png", Payload_Contractor_Bar_Percentage)
            pic_left_3  = int(prs.slide_width *0.24)
            pic_top_3   = int(prs.slide_width *0.4)
            pic_width_3 = int(prs.slide_width *0.47)


            pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Payload_"+Contractor+"_Bar.png", pic_left_1, pic_top_1, pic_width_1)
            pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\Payload_"+Contractor+"_Bar_Percentage.png", pic_left_3, pic_top_3, pic_width_3)

            prs.save('test.pptx')


# ******************************************************************************************
# (((((((((((((((((((((((((((((((((((((   CC and RD    )))))))))))))))))))))))))))))))))))))
# ******************************************************************************************

conn_performanceDB.execute("select Wk,Contractor, avg([CC2]) as 'CC2', avg([CC3]) as 'CC3', avg([CC]) as 'CC', sum([2G_TCH_Traffic]) as '2G_TCH_Traffic', sum([3G_CS_Traffic]) as '3G_CS_Traffic', sum([4G_Volte_Traffic]) as '4G_Volte_Traffic',  sum([Total_Traffic]) as 'Total_Traffic' from ("+
                           "select Wk, Contractor,PIndex, SUM([2G TCH Traffic]*[CC2 (%)])/sum([2G TCH Traffic]) as 'CC2',"+
                              "SUM([3G_CS_Traffic]*[CC3 (%)])/sum([3G_CS_Traffic]) as 'CC3',"+
							  "SUM(([2G TCH Traffic]+[3G_CS_Traffic])*[CC (%)])/sum([2G TCH Traffic]+[3G_CS_Traffic]) as 'CC',"+
							  "avg([2G TCH Traffic]) as '2G_TCH_Traffic',"+
							  "avg([3G_CS_Traffic]) as '3G_CS_Traffic',"+
							  "avg([L1800 Traffic]+[L2600 Traffic]+[L2300 Traffic]+[L2100 Traffic]) as '4G_Volte_Traffic',"+
							  "avg([Total Voice Traffic (Erlang)]) as 'Total_Traffic' "+
							  "from  Province_KPI_Score_Band_CS_Daily  group by Wk, Contractor, PIndex ) tble  group by Wk, Contractor  order by Wk")
CC_Table=conn_performanceDB.fetchall()


conn_performanceDB.execute("select Wk,Contractor, avg([RD2]) as 'RD2', avg([RD3]) as 'RD3', avg([RD4]) as 'RD4', avg([RD]) as 'RD', sum([2G_PS_Payload]) as '2G_PS_Payload', sum([3G_PS_Payload]) as '3G_PS_Payload', sum([4G_PS_Payload]) as '4G_PS_Payload', sum([Total_Payload]) as 'Total_Payload' from ("+
                             "select Wk, Contractor,[Province Index], SUM([2G PS Traffic (GB)]*[RD2 (%)])/sum([2G PS Traffic (GB)]) as 'RD2',"+
                              "SUM([3G Payload (GB)]*[RD3 (%)])/sum([3G Payload (GB)]) as 'RD3',"+
							  "SUM([4G Payload (GB)]*[RD4 (%)])/sum([4G Payload (GB)]) as 'RD4',"+
							  "SUM([Total Payload (GB)]*[RD (%)])/sum([Total Payload (GB)]) as 'RD',"+
							  "avg([2G PS Traffic (GB)]) as '2G_PS_Payload',"+
							  "avg([3G Payload (GB)]) as '3G_PS_Payload',"+
							  "avg([4G Payload (GB)]) as '4G_PS_Payload',"+
							  "avg([Total Payload (GB)]) as 'Total_Payload' "+
							  "from  Province_KPI_Score_Band_PS_Daily group by Wk, Contractor, [Province Index] ) tble group by Wk, Contractor  order by Wk")
RD_Table=conn_performanceDB.fetchall()


CC2_NAK_Alborz=[]
CC2_NAK_Tehran=[]
CC2_NAK_North=[]
CC2_NAK_Nokia=[]
CC2_NAK_Huawei=[]
CC2_Farafan=[]
CC2_BR_TEL=[]
CC2_Huawei=[]
CC2_Iran=[]

CC3_NAK_Alborz=[]
CC3_NAK_Tehran=[]
CC3_NAK_North=[]
CC3_NAK_Nokia=[]
CC3_NAK_Huawei=[]
CC3_Farafan=[]
CC3_BR_TEL=[]
CC3_Huawei=[]
CC3_Iran=[]

CC_NAK_Alborz=[]
CC_NAK_Tehran=[]
CC_NAK_North=[]
CC_NAK_Nokia=[]
CC_NAK_Huawei=[]
CC_Farafan=[]
CC_BR_TEL=[]
CC_Huawei=[]
CC_Iran=[]



CS_Traffic_2G_NAK_Alborz=[]
CS_Traffic_2G_NAK_Tehran=[]
CS_Traffic_2G_NAK_North=[]
CS_Traffic_2G_NAK_Nokia=[]
CS_Traffic_2G_NAK_Huawei=[]
CS_Traffic_2G_Farafan=[]
CS_Traffic_2G_BR_TEL=[]
CS_Traffic_2G_Huawei=[]
CS_Traffic_2G_Iran=[]

CS_Traffic_3G_NAK_Alborz=[]
CS_Traffic_3G_NAK_Tehran=[]
CS_Traffic_3G_NAK_North=[]
CS_Traffic_3G_NAK_Nokia=[]
CS_Traffic_3G_NAK_Huawei=[]
CS_Traffic_3G_Farafan=[]
CS_Traffic_3G_BR_TEL=[]
CS_Traffic_3G_Huawei=[]
CS_Traffic_3G_Iran=[]

CS_Traffic_4G_NAK_Alborz=[]
CS_Traffic_4G_NAK_Tehran=[]
CS_Traffic_4G_NAK_North=[]
CS_Traffic_4G_NAK_Nokia=[]
CS_Traffic_4G_NAK_Huawei=[]
CS_Traffic_4G_Farafan=[]
CS_Traffic_4G_BR_TEL=[]
CS_Traffic_4G_Huawei=[]
CS_Traffic_4G_Iran=[]

Week_Vec=[]

for i in range(len(CC_Table)):
    Row_Data=str(CC_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Contractor=Row_Data[1]
    CC2_Str=Row_Data[2]
    CC2_Val=round(float(CC2_Str[0:len(CC2_Str)-1]),2)
    CC3_Str=Row_Data[3]
    CC3_Val=round(float(CC3_Str[0:len(CC3_Str)-1]),2)
    CC_Str=Row_Data[4]
    CC_Val=round(float(CC_Str[0:len(CC_Str)-1]),2)
    CS_2G_Str=Row_Data[5]
    CS_2G_Val=round(float(CS_2G_Str[0:len(CS_2G_Str)-1]),2)
    CS_3G_Str=Row_Data[6]
    CS_3G_Val=round(float(CS_3G_Str[0:len(CS_3G_Str)-1]),2)
    CS_4G_Str=Row_Data[7]
    CS_4G_Val=round(float(CS_4G_Str[0:len(CS_4G_Str)-1]),2)
    CS_Str=Row_Data[8]
    CS_Val=round(float(CS_Str[0:len(CS_Str)-1]),2)

    Week=Week[2:9]
    Contractor=Contractor[1:len(Contractor)-1]
    

    #if Week=='1401-33':
    #    break

    if (Contractor=='NAK-Alborz'):
        Week_Vec.append('W'+Week[5:7])
        CC2_NAK_Alborz.append(CC2_Val)
        CC3_NAK_Alborz.append(CC3_Val)
        CC_NAK_Alborz.append(CC_Val)
        CS_Traffic_2G_NAK_Alborz.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Alborz.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Alborz.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-Tehran'):
        CC2_NAK_Tehran.append(CC2_Val)
        CC3_NAK_Tehran.append(CC3_Val)
        CC_NAK_Tehran.append(CC_Val)
        CS_Traffic_2G_NAK_Tehran.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Tehran.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Tehran.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-North'):
        CC2_NAK_North.append(CC2_Val)
        CC3_NAK_North.append(CC3_Val)
        CC_NAK_North.append(CC_Val)
        CS_Traffic_2G_NAK_North.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_North.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_North.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-Nokia'):
        CC2_NAK_Nokia.append(CC2_Val)
        CC3_NAK_Nokia.append(CC3_Val)
        CC_NAK_Nokia.append(CC_Val)
        CS_Traffic_2G_NAK_Nokia.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Nokia.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Nokia.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-Huawei'):
        CC2_NAK_Huawei.append(CC2_Val)
        CC3_NAK_Huawei.append(CC3_Val)
        CC_NAK_Huawei.append(CC_Val)
        CS_Traffic_2G_NAK_Huawei.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Huawei.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Huawei.append(CS_4G_Val/1e3)
    if (Contractor=='Farafan'):
        CC2_Farafan.append(CC2_Val)
        CC3_Farafan.append(CC3_Val)
        CC_Farafan.append(CC_Val)
        CS_Traffic_2G_Farafan.append(CS_2G_Val/1e3)
        CS_Traffic_3G_Farafan.append(CS_3G_Val/1e3)
        CS_Traffic_4G_Farafan.append(CS_4G_Val/1e3)
    if (Contractor=='BR-TEL'):
        CC2_BR_TEL.append(CC2_Val)
        CC3_BR_TEL.append(CC3_Val)
        CC_BR_TEL.append(CC_Val)
        CS_Traffic_2G_BR_TEL.append(CS_2G_Val/1e3)
        CS_Traffic_3G_BR_TEL.append(CS_3G_Val/1e3)
        CS_Traffic_4G_BR_TEL.append(CS_4G_Val/1e3)
    if (Contractor=='Huawei'):
        CC2_Huawei.append(CC2_Val)
        CC3_Huawei.append(CC3_Val)
        CC_Huawei.append(CC_Val)
        CS_Traffic_2G_Huawei.append(CS_2G_Val/1e3)
        CS_Traffic_3G_Huawei.append(CS_3G_Val/1e3)
        CS_Traffic_4G_Huawei.append(CS_4G_Val/1e3)
    if (Contractor=='IRAN'):
        CC2_Iran.append(CC2_Val)
        CC3_Iran.append(CC3_Val)
        CC_Iran.append(CC_Val)
        CS_Traffic_2G_Iran.append(CS_2G_Val/1e3)
        CS_Traffic_3G_Iran.append(CS_3G_Val/1e3)
        CS_Traffic_4G_Iran.append(CS_4G_Val/1e3)



Last_CC2=[CC2_NAK_Alborz[len(CC2_NAK_Alborz)-1], CC2_NAK_Tehran[len(CC2_NAK_Alborz)-1], CC2_NAK_North[len(CC2_NAK_Alborz)-1], CC2_NAK_Nokia[len(CC2_NAK_Alborz)-1], CC2_NAK_Huawei[len(CC2_NAK_Alborz)-1], CC2_Farafan[len(CC2_NAK_Alborz)-1], CC2_BR_TEL[len(CC2_NAK_Alborz)-1], CC2_Huawei[len(CC2_NAK_Alborz)-1] ]
Index_of_Sort_CC2=np.argsort(Last_CC2)
Data_Sorted_Array_CC2=[]
x_Labels_CC2=[];
for k in range(len(Index_of_Sort_CC2)):
    if Index_of_Sort_CC2[k]==0:
        Data_Sorted_Array_CC2.append(CC2_NAK_Alborz)
        x_Labels_CC2.append('NAK-Alborz')
    if Index_of_Sort_CC2[k]==1:
        Data_Sorted_Array_CC2.append(CC2_NAK_Tehran)
        x_Labels_CC2.append('NAK-Tehran')
    if Index_of_Sort_CC2[k]==2:
        Data_Sorted_Array_CC2.append(CC2_NAK_North)
        x_Labels_CC2.append('NAK-North')
    if Index_of_Sort_CC2[k]==3:
        Data_Sorted_Array_CC2.append(CC2_NAK_Nokia)
        x_Labels_CC2.append('NAK-Nokia')
    if Index_of_Sort_CC2[k]==4:
        Data_Sorted_Array_CC2.append(CC2_NAK_Huawei)
        x_Labels_CC2.append('NAK-Huawei')
    if Index_of_Sort_CC2[k]==5:
        Data_Sorted_Array_CC2.append(CC2_Farafan)
        x_Labels_CC2.append('Farafan')
    if Index_of_Sort_CC2[k]==6:
        Data_Sorted_Array_CC2.append(CC2_BR_TEL)
        x_Labels_CC2.append('BR-TEL')
    if Index_of_Sort_CC2[k]==7:
        Data_Sorted_Array_CC2.append(CC2_Huawei)
        x_Labels_CC2.append('Huawei')


data=np.array(Data_Sorted_Array_CC2)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_CC2):
    plt.text( i + dx[31],Last_CC2[Index_of_Sort_CC2[i]] , str(Last_CC2[Index_of_Sort_CC2[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_CC2)
font1 = {'family':'serif','color':'black','size':17}
plt.title("CC2(%)", fontdict = font1)
plt.ylim(75, 100)
grid(True)
plt.savefig('CC2.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC2.png")
y=80
x=10
h=1000
w=520
CC2_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("CC2.png", CC2_Bar_Cropped)



Last_CC3=[CC3_NAK_Alborz[len(CC3_NAK_Alborz)-1], CC3_NAK_Tehran[len(CC3_NAK_Alborz)-1], CC3_NAK_North[len(CC3_NAK_Alborz)-1], CC3_NAK_Nokia[len(CC3_NAK_Alborz)-1], CC3_NAK_Huawei[len(CC3_NAK_Alborz)-1], CC3_Farafan[len(CC3_NAK_Alborz)-1], CC3_BR_TEL[len(CC3_NAK_Alborz)-1], CC3_Huawei[len(CC3_NAK_Alborz)-1] ]
Index_of_Sort_CC3=np.argsort(Last_CC3)
Data_Sorted_Array_CC3=[]
x_Labels_CC3=[];
for k in range(len(Index_of_Sort_CC3)):
    if Index_of_Sort_CC3[k]==0:
        Data_Sorted_Array_CC3.append(CC3_NAK_Alborz)
        x_Labels_CC3.append('NAK-Alborz')
    if Index_of_Sort_CC3[k]==1:
        Data_Sorted_Array_CC3.append(CC3_NAK_Tehran)
        x_Labels_CC3.append('NAK-Tehran')
    if Index_of_Sort_CC3[k]==2:
        Data_Sorted_Array_CC3.append(CC3_NAK_North)
        x_Labels_CC3.append('NAK-North')
    if Index_of_Sort_CC3[k]==3:
        Data_Sorted_Array_CC3.append(CC3_NAK_Nokia)
        x_Labels_CC3.append('NAK-Nokia')
    if Index_of_Sort_CC3[k]==4:
        Data_Sorted_Array_CC3.append(CC3_NAK_Huawei)
        x_Labels_CC3.append('NAK-Huawei')
    if Index_of_Sort_CC3[k]==5:
        Data_Sorted_Array_CC3.append(CC3_Farafan)
        x_Labels_CC3.append('Farafan')
    if Index_of_Sort_CC3[k]==6:
        Data_Sorted_Array_CC3.append(CC3_BR_TEL)
        x_Labels_CC3.append('BR-TEL')
    if Index_of_Sort_CC3[k]==7:
        Data_Sorted_Array_CC3.append(CC3_Huawei)
        x_Labels_CC3.append('Huawei')


data=np.array(Data_Sorted_Array_CC3)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_CC3):
    plt.text( i + dx[31],Last_CC3[Index_of_Sort_CC3[i]] , str(Last_CC3[Index_of_Sort_CC3[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_CC3)
font1 = {'family':'serif','color':'black','size':17}
plt.title("CC3(%)", fontdict = font1)
plt.ylim(75, 100)
grid(True)
plt.savefig('CC3.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC3.png")
y=80
x=10
h=1000
w=520
CC3_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("CC3.png", CC3_Bar_Cropped)

Last_CC=[CC_NAK_Alborz[len(CC_NAK_Alborz)-1], CC_NAK_Tehran[len(CC_NAK_Alborz)-1], CC_NAK_North[len(CC_NAK_Alborz)-1], CC_NAK_Nokia[len(CC_NAK_Alborz)-1], CC_NAK_Huawei[len(CC_NAK_Alborz)-1], CC_Farafan[len(CC_NAK_Alborz)-1], CC_BR_TEL[len(CC_NAK_Alborz)-1], CC_Huawei[len(CC_NAK_Alborz)-1] ]
Index_of_Sort_CC=np.argsort(Last_CC)
Data_Sorted_Array_CC=[]
x_Labels_CC=[];
for k in range(len(Index_of_Sort_CC)):
    if Index_of_Sort_CC[k]==0:
        Data_Sorted_Array_CC.append(CC_NAK_Alborz)
        x_Labels_CC.append('NAK-Alborz')
    if Index_of_Sort_CC[k]==1:
        Data_Sorted_Array_CC.append(CC_NAK_Tehran)
        x_Labels_CC.append('NAK-Tehran')
    if Index_of_Sort_CC[k]==2:
        Data_Sorted_Array_CC.append(CC_NAK_North)
        x_Labels_CC.append('NAK-North')
    if Index_of_Sort_CC[k]==3:
        Data_Sorted_Array_CC.append(CC_NAK_Nokia)
        x_Labels_CC.append('NAK-Nokia')
    if Index_of_Sort_CC[k]==4:
        Data_Sorted_Array_CC.append(CC_NAK_Huawei)
        x_Labels_CC.append('NAK-Huawei')
    if Index_of_Sort_CC[k]==5:
        Data_Sorted_Array_CC.append(CC_Farafan)
        x_Labels_CC.append('Farafan')
    if Index_of_Sort_CC[k]==6:
        Data_Sorted_Array_CC.append(CC_BR_TEL)
        x_Labels_CC.append('BR-TEL')
    if Index_of_Sort_CC[k]==7:
        Data_Sorted_Array_CC.append(CC_Huawei)
        x_Labels_CC.append('Huawei')


data=np.array(Data_Sorted_Array_CC)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_CC):
    plt.text( i + dx[31],Last_CC[Index_of_Sort_CC[i]] , str(Last_CC[Index_of_Sort_CC[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_CC)
font1 = {'family':'serif','color':'black','size':17}
plt.title("CC(%)", fontdict = font1)
plt.ylim(75, 100)
grid(True)
plt.savefig('CC.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC.png")
y=80
x=10
h=1000
w=520
CC_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("CC.png", CC_Bar_Cropped)




RD2_NAK_Alborz=[]
RD2_NAK_Tehran=[]
RD2_NAK_North=[]
RD2_NAK_Nokia=[]
RD2_NAK_Huawei=[]
RD2_Farafan=[]
RD2_BR_TEL=[]
RD2_Huawei=[]
RD2_Iran=[]

RD3_NAK_Alborz=[]
RD3_NAK_Tehran=[]
RD3_NAK_North=[]
RD3_NAK_Nokia=[]
RD3_NAK_Huawei=[]
RD3_Farafan=[]
RD3_BR_TEL=[]
RD3_Huawei=[]
RD3_Iran=[]

RD4_NAK_Alborz=[]
RD4_NAK_Tehran=[]
RD4_NAK_North=[]
RD4_NAK_Nokia=[]
RD4_NAK_Huawei=[]
RD4_Farafan=[]
RD4_BR_TEL=[]
RD4_Huawei=[]
RD4_Iran=[]

RD_NAK_Alborz=[]
RD_NAK_Tehran=[]
RD_NAK_North=[]
RD_NAK_Nokia=[]
RD_NAK_Huawei=[]
RD_Farafan=[]
RD_BR_TEL=[]
RD_Huawei=[]
RD_Iran=[]


PS_Traffic_2G_NAK_Alborz=[]
PS_Traffic_2G_NAK_Tehran=[]
PS_Traffic_2G_NAK_North=[]
PS_Traffic_2G_NAK_Nokia=[]
PS_Traffic_2G_NAK_Huawei=[]
PS_Traffic_2G_Farafan=[]
PS_Traffic_2G_BR_TEL=[]
PS_Traffic_2G_Huawei=[]
PS_Traffic_2G_Iran=[]

PS_Traffic_3G_NAK_Alborz=[]
PS_Traffic_3G_NAK_Tehran=[]
PS_Traffic_3G_NAK_North=[]
PS_Traffic_3G_NAK_Nokia=[]
PS_Traffic_3G_NAK_Huawei=[]
PS_Traffic_3G_Farafan=[]
PS_Traffic_3G_BR_TEL=[]
PS_Traffic_3G_Huawei=[]
PS_Traffic_3G_Iran=[]

PS_Traffic_4G_NAK_Alborz=[]
PS_Traffic_4G_NAK_Tehran=[]
PS_Traffic_4G_NAK_North=[]
PS_Traffic_4G_NAK_Nokia=[]
PS_Traffic_4G_NAK_Huawei=[]
PS_Traffic_4G_Farafan=[]
PS_Traffic_4G_BR_TEL=[]
PS_Traffic_4G_Huawei=[]
PS_Traffic_4G_Iran=[]



for i in range(len(RD_Table)):
    Row_Data=str(RD_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Contractor=Row_Data[1]
    RD2_Str=Row_Data[2]
    RD2_Val=round(float(RD2_Str[0:len(RD2_Str)-1]),2)
    RD3_Str=Row_Data[3]
    RD3_Val=round(float(RD3_Str[0:len(RD3_Str)-1]),2)
    RD4_Str=Row_Data[4]
    RD4_Val=round(float(RD4_Str[0:len(RD4_Str)-1]),2)
    RD_Str=Row_Data[5]
    RD_Val=round(float(RD_Str[0:len(RD_Str)-1]),2)
    PS_2G_Str=Row_Data[6]
    PS_2G_Val=round(float(PS_2G_Str[0:len(PS_2G_Str)-1]),2)
    PS_3G_Str=Row_Data[7]
    PS_3G_Val=round(float(PS_3G_Str[0:len(PS_3G_Str)-1]),2)
    PS_4G_Str=Row_Data[8]
    PS_4G_Val=round(float(PS_4G_Str[0:len(PS_4G_Str)-1]),2)
    PS_Str=Row_Data[9]
    PS_Val=round(float(PS_Str[0:len(PS_Str)-1]),2)
    Week=Week[2:9]
    Contractor=Contractor[1:len(Contractor)-1]
    

    #if Week=='1401-33':
    #    break

    if (Contractor=='NAK-Alborz'):
        RD2_NAK_Alborz.append(RD2_Val)
        RD3_NAK_Alborz.append(RD3_Val)
        RD4_NAK_Alborz.append(RD4_Val)
        RD_NAK_Alborz.append(RD_Val)
        PS_Traffic_2G_NAK_Alborz.append(PS_2G_Val/1e3)
        PS_Traffic_3G_NAK_Alborz.append(PS_3G_Val/1e3)
        PS_Traffic_4G_NAK_Alborz.append(PS_4G_Val/1e3)
    if (Contractor=='NAK-Tehran'):
        RD2_NAK_Tehran.append(RD2_Val)
        RD3_NAK_Tehran.append(RD3_Val)
        RD4_NAK_Tehran.append(RD4_Val)
        RD_NAK_Tehran.append(RD_Val)
        PS_Traffic_2G_NAK_Tehran.append(PS_2G_Val/1e3)
        PS_Traffic_3G_NAK_Tehran.append(PS_3G_Val/1e3)
        PS_Traffic_4G_NAK_Tehran.append(PS_4G_Val/1e3)
    if (Contractor=='NAK-North'):
        RD2_NAK_North.append(RD2_Val)
        RD3_NAK_North.append(RD3_Val)
        RD4_NAK_North.append(RD4_Val)
        RD_NAK_North.append(RD_Val)
        PS_Traffic_2G_NAK_North.append(PS_2G_Val/1e3)
        PS_Traffic_3G_NAK_North.append(PS_3G_Val/1e3)
        PS_Traffic_4G_NAK_North.append(PS_4G_Val/1e3)
    if (Contractor=='NAK-Nokia'):
        RD2_NAK_Nokia.append(RD2_Val)
        RD3_NAK_Nokia.append(RD3_Val)
        RD4_NAK_Nokia.append(RD4_Val)
        RD_NAK_Nokia.append(RD_Val)
        PS_Traffic_2G_NAK_Nokia.append(PS_2G_Val/1e3)
        PS_Traffic_3G_NAK_Nokia.append(PS_3G_Val/1e3)
        PS_Traffic_4G_NAK_Nokia.append(PS_4G_Val/1e3)
    if (Contractor=='NAK-Huawei'):
        RD2_NAK_Huawei.append(RD2_Val)
        RD3_NAK_Huawei.append(RD3_Val)
        RD4_NAK_Huawei.append(RD4_Val)
        RD_NAK_Huawei.append(RD_Val)
        PS_Traffic_2G_NAK_Huawei.append(PS_2G_Val/1e3)
        PS_Traffic_3G_NAK_Huawei.append(PS_3G_Val/1e3)
        PS_Traffic_4G_NAK_Huawei.append(PS_4G_Val/1e3)
    if (Contractor=='Farafan'):
        RD2_Farafan.append(RD2_Val)
        RD3_Farafan.append(RD3_Val)
        RD4_Farafan.append(RD4_Val)
        RD_Farafan.append(RD_Val)
        PS_Traffic_2G_Farafan.append(PS_2G_Val/1e3)
        PS_Traffic_3G_Farafan.append(PS_3G_Val/1e3)
        PS_Traffic_4G_Farafan.append(PS_4G_Val/1e3)
    if (Contractor=='BR-TEL'):
        RD2_BR_TEL.append(RD2_Val)
        RD3_BR_TEL.append(RD3_Val)
        RD4_BR_TEL.append(RD4_Val)
        RD_BR_TEL.append(RD_Val)
        PS_Traffic_2G_BR_TEL.append(PS_2G_Val/1e3)
        PS_Traffic_3G_BR_TEL.append(PS_3G_Val/1e3)
        PS_Traffic_4G_BR_TEL.append(PS_4G_Val/1e3)
    if (Contractor=='Huawei'):
        RD2_Huawei.append(RD2_Val)
        RD3_Huawei.append(RD3_Val)
        RD4_Huawei.append(RD4_Val)
        RD_Huawei.append(RD_Val)
        PS_Traffic_2G_Huawei.append(PS_2G_Val/1e3)
        PS_Traffic_3G_Huawei.append(PS_3G_Val/1e3)
        PS_Traffic_4G_Huawei.append(PS_4G_Val/1e3)
    if (Contractor=='IRAN'):
        RD2_Iran.append(RD2_Val)
        RD3_Iran.append(RD3_Val)
        RD4_Iran.append(RD4_Val)
        RD_Iran.append(RD_Val)
        PS_Traffic_2G_Iran.append(PS_2G_Val/1e3)
        PS_Traffic_3G_Iran.append(PS_3G_Val/1e3)
        PS_Traffic_4G_Iran.append(PS_4G_Val/1e3)




Last_RD2=[RD2_NAK_Alborz[len(RD2_NAK_Alborz)-1], RD2_NAK_Tehran[len(RD2_NAK_Alborz)-1], RD2_NAK_North[len(RD2_NAK_Alborz)-1], RD2_NAK_Nokia[len(RD2_NAK_Alborz)-1], RD2_NAK_Huawei[len(RD2_NAK_Alborz)-1], RD2_Farafan[len(RD2_NAK_Alborz)-1], RD2_BR_TEL[len(RD2_NAK_Alborz)-1], RD2_Huawei[len(RD2_NAK_Alborz)-1] ]
Index_of_Sort_RD2=np.argsort(Last_RD2)
Data_Sorted_Array_RD2=[]
x_Labels_RD2=[];
for k in range(len(Index_of_Sort_RD2)):
    if Index_of_Sort_RD2[k]==0:
        Data_Sorted_Array_RD2.append(RD2_NAK_Alborz)
        x_Labels_RD2.append('NAK-Alborz')
    if Index_of_Sort_RD2[k]==1:
        Data_Sorted_Array_RD2.append(RD2_NAK_Tehran)
        x_Labels_RD2.append('NAK-Tehran')
    if Index_of_Sort_RD2[k]==2:
        Data_Sorted_Array_RD2.append(RD2_NAK_North)
        x_Labels_RD2.append('NAK-North')
    if Index_of_Sort_RD2[k]==3:
        Data_Sorted_Array_RD2.append(RD2_NAK_Nokia)
        x_Labels_RD2.append('NAK-Nokia')
    if Index_of_Sort_RD2[k]==4:
        Data_Sorted_Array_RD2.append(RD2_NAK_Huawei)
        x_Labels_RD2.append('NAK-Huawei')
    if Index_of_Sort_RD2[k]==5:
        Data_Sorted_Array_RD2.append(RD2_Farafan)
        x_Labels_RD2.append('Farafan')
    if Index_of_Sort_RD2[k]==6:
        Data_Sorted_Array_RD2.append(RD2_BR_TEL)
        x_Labels_RD2.append('BR-TEL')
    if Index_of_Sort_RD2[k]==7:
        Data_Sorted_Array_RD2.append(RD2_Huawei)
        x_Labels_RD2.append('Huawei')


data=np.array(Data_Sorted_Array_RD2)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(6)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_RD2):
    plt.text( i + dx[31],Last_RD2[Index_of_Sort_RD2[i]] , str(Last_RD2[Index_of_Sort_RD2[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_RD2)
font1 = {'family':'serif','color':'black','size':14}
plt.title("RD2(%)", fontdict = font1)
plt.ylim(40, 100)
grid(True)
plt.savefig('RD2.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD2.png")
y=80
x=5
h=1000
w=520
RD2_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("RD2.png", RD2_Bar_Cropped)



Last_RD3=[RD3_NAK_Alborz[len(RD3_NAK_Alborz)-1], RD3_NAK_Tehran[len(RD3_NAK_Alborz)-1], RD3_NAK_North[len(RD3_NAK_Alborz)-1], RD3_NAK_Nokia[len(RD3_NAK_Alborz)-1], RD3_NAK_Huawei[len(RD3_NAK_Alborz)-1], RD3_Farafan[len(RD3_NAK_Alborz)-1], RD3_BR_TEL[len(RD3_NAK_Alborz)-1], RD3_Huawei[len(RD3_NAK_Alborz)-1] ]
Index_of_Sort_RD3=np.argsort(Last_RD3)
Data_Sorted_Array_RD3=[]
x_Labels_RD3=[];
for k in range(len(Index_of_Sort_RD3)):
    if Index_of_Sort_RD3[k]==0:
        Data_Sorted_Array_RD3.append(RD3_NAK_Alborz)
        x_Labels_RD3.append('NAK-Alborz')
    if Index_of_Sort_RD3[k]==1:
        Data_Sorted_Array_RD3.append(RD3_NAK_Tehran)
        x_Labels_RD3.append('NAK-Tehran')
    if Index_of_Sort_RD3[k]==2:
        Data_Sorted_Array_RD3.append(RD3_NAK_North)
        x_Labels_RD3.append('NAK-North')
    if Index_of_Sort_RD3[k]==3:
        Data_Sorted_Array_RD3.append(RD3_NAK_Nokia)
        x_Labels_RD3.append('NAK-Nokia')
    if Index_of_Sort_RD3[k]==4:
        Data_Sorted_Array_RD3.append(RD3_NAK_Huawei)
        x_Labels_RD3.append('NAK-Huawei')
    if Index_of_Sort_RD3[k]==5:
        Data_Sorted_Array_RD3.append(RD3_Farafan)
        x_Labels_RD3.append('Farafan')
    if Index_of_Sort_RD3[k]==6:
        Data_Sorted_Array_RD3.append(RD3_BR_TEL)
        x_Labels_RD3.append('BR-TEL')
    if Index_of_Sort_RD3[k]==7:
        Data_Sorted_Array_RD3.append(RD3_Huawei)
        x_Labels_RD3.append('Huawei')


data=np.array(Data_Sorted_Array_RD3)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(6)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_RD3):
    plt.text( i + dx[31],Last_RD3[Index_of_Sort_RD3[i]] , str(Last_RD3[Index_of_Sort_RD3[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_RD3)
font1 = {'family':'serif','color':'black','size':14}
plt.title("RD3(%)", fontdict = font1)
plt.ylim(70, 100)
grid(True)
plt.savefig('RD3.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD3.png")
y=80
x=5
h=1000
w=520
RD3_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("RD3.png", RD3_Bar_Cropped)



Last_RD4=[RD4_NAK_Alborz[len(RD4_NAK_Alborz)-1], RD4_NAK_Tehran[len(RD4_NAK_Alborz)-1], RD4_NAK_North[len(RD4_NAK_Alborz)-1], RD4_NAK_Nokia[len(RD4_NAK_Alborz)-1], RD4_NAK_Huawei[len(RD4_NAK_Alborz)-1], RD4_Farafan[len(RD4_NAK_Alborz)-1], RD4_BR_TEL[len(RD4_NAK_Alborz)-1], RD4_Huawei[len(RD4_NAK_Alborz)-1] ]
Index_of_Sort_RD4=np.argsort(Last_RD4)
Data_Sorted_Array_RD4=[]
x_Labels_RD4=[];
for k in range(len(Index_of_Sort_RD4)):
    if Index_of_Sort_RD4[k]==0:
        Data_Sorted_Array_RD4.append(RD4_NAK_Alborz)
        x_Labels_RD4.append('NAK-Alborz')
    if Index_of_Sort_RD4[k]==1:
        Data_Sorted_Array_RD4.append(RD4_NAK_Tehran)
        x_Labels_RD4.append('NAK-Tehran')
    if Index_of_Sort_RD4[k]==2:
        Data_Sorted_Array_RD4.append(RD4_NAK_North)
        x_Labels_RD4.append('NAK-North')
    if Index_of_Sort_RD4[k]==3:
        Data_Sorted_Array_RD4.append(RD4_NAK_Nokia)
        x_Labels_RD4.append('NAK-Nokia')
    if Index_of_Sort_RD4[k]==4:
        Data_Sorted_Array_RD4.append(RD4_NAK_Huawei)
        x_Labels_RD4.append('NAK-Huawei')
    if Index_of_Sort_RD4[k]==5:
        Data_Sorted_Array_RD4.append(RD4_Farafan)
        x_Labels_RD4.append('Farafan')
    if Index_of_Sort_RD4[k]==6:
        Data_Sorted_Array_RD4.append(RD4_BR_TEL)
        x_Labels_RD4.append('BR-TEL')
    if Index_of_Sort_RD4[k]==7:
        Data_Sorted_Array_RD4.append(RD4_Huawei)
        x_Labels_RD4.append('Huawei')


data=np.array(Data_Sorted_Array_RD4)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(6)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_RD4):
    plt.text( i + dx[31],Last_RD4[Index_of_Sort_RD4[i]] , str(Last_RD4[Index_of_Sort_RD4[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_RD4)
font1 = {'family':'serif','color':'black','size':14}
plt.title("RD4(%)", fontdict = font1)
plt.ylim(70, 100)
grid(True)
plt.savefig('RD4.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD4.png")
y=80
x=5
h=1000
w=520
RD4_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("RD4.png", RD4_Bar_Cropped)



Last_RD=[RD_NAK_Alborz[len(RD_NAK_Alborz)-1], RD_NAK_Tehran[len(RD_NAK_Alborz)-1], RD_NAK_North[len(RD_NAK_Alborz)-1], RD_NAK_Nokia[len(RD_NAK_Alborz)-1], RD_NAK_Huawei[len(RD_NAK_Alborz)-1], RD_Farafan[len(RD_NAK_Alborz)-1], RD_BR_TEL[len(RD_NAK_Alborz)-1], RD_Huawei[len(RD_NAK_Alborz)-1] ]
Index_of_Sort_RD=np.argsort(Last_RD)
Data_Sorted_Array_RD=[]
x_Labels_RD=[];
for k in range(len(Index_of_Sort_RD)):
    if Index_of_Sort_RD[k]==0:
        Data_Sorted_Array_RD.append(RD_NAK_Alborz)
        x_Labels_RD.append('NAK-Alborz')
    if Index_of_Sort_RD[k]==1:
        Data_Sorted_Array_RD.append(RD_NAK_Tehran)
        x_Labels_RD.append('NAK-Tehran')
    if Index_of_Sort_RD[k]==2:
        Data_Sorted_Array_RD.append(RD_NAK_North)
        x_Labels_RD.append('NAK-North')
    if Index_of_Sort_RD[k]==3:
        Data_Sorted_Array_RD.append(RD_NAK_Nokia)
        x_Labels_RD.append('NAK-Nokia')
    if Index_of_Sort_RD[k]==4:
        Data_Sorted_Array_RD.append(RD_NAK_Huawei)
        x_Labels_RD.append('NAK-Huawei')
    if Index_of_Sort_RD[k]==5:
        Data_Sorted_Array_RD.append(RD_Farafan)
        x_Labels_RD.append('Farafan')
    if Index_of_Sort_RD[k]==6:
        Data_Sorted_Array_RD.append(RD_BR_TEL)
        x_Labels_RD.append('BR-TEL')
    if Index_of_Sort_RD[k]==7:
        Data_Sorted_Array_RD.append(RD_Huawei)
        x_Labels_RD.append('Huawei')


data=np.array(Data_Sorted_Array_RD)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "green",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_RD):
    plt.text( i + dx[31],Last_RD[Index_of_Sort_RD[i]] , str(Last_RD[Index_of_Sort_RD[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_RD)
font1 = {'family':'serif','color':'black','size':17}
plt.title("RD(%)", fontdict = font1)
plt.ylim(70, 100)
grid(True)
plt.savefig('RD.png')



image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD.png")
y=80
x=10
h=1000
w=520
RD_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("RD.png", RD_Bar_Cropped)





pic_left_1  = int(prs.slide_width *0.047)
pic_top_1   = int(prs.slide_width *0.03)
pic_width_1 = int(prs.slide_width *0.9)


pic_left_3  = int(prs.slide_width *0.047)
pic_top_3   = int(prs.slide_width *0.38)
pic_width_3 = int(prs.slide_width *0.9)

slide = prs.slides.add_slide(blank_slide_layout)

pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC.png", pic_left_1, pic_top_1, pic_width_1)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD.png", pic_left_3, pic_top_3, pic_width_3)

slide = prs.slides.add_slide(blank_slide_layout)

pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC2.png", pic_left_1, pic_top_1, pic_width_1)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC3.png", pic_left_3, pic_top_3, pic_width_3)


pic_left_1  = int(prs.slide_width *0.047)
pic_top_1   = int(prs.slide_width *0.03)
pic_width_1 = int(prs.slide_width *0.9)

pic_left_2  = int(prs.slide_width *0.047)
pic_top_2   = int(prs.slide_width *0.27)
pic_width_2 = int(prs.slide_width *0.9)

pic_left_3  = int(prs.slide_width *0.047)
pic_top_3   = int(prs.slide_width *0.51)
pic_width_3 = int(prs.slide_width *0.9)



slide = prs.slides.add_slide(blank_slide_layout)

pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD2.png", pic_left_1, pic_top_1, pic_width_1)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD3.png", pic_left_2, pic_top_2, pic_width_2)
pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD4.png", pic_left_3, pic_top_3, pic_width_3)


prs.save('test.pptx')




# *********************************************************************************************************
# (((((((((((((((((((((((((((((((((((((   CC and RD per Contractor    )))))))))))))))))))))))))))))))))))))
# *********************************************************************************************************

for t in range(8):
        if (t==0 ):
            Contractor="NAK-Alborz"
            CS_Traffic=CS_Traffic_NAK_Alborz
            CC=CC_NAK_Alborz
            PS_Traffic=PS_Traffic_NAK_Alborz
            RD=RD_NAK_Alborz
            CS_Traffic_2G=CS_Traffic_2G_NAK_Alborz
            CC2=CC2_NAK_Alborz
            CS_Traffic_3G=CS_Traffic_3G_NAK_Alborz
            CC3=CC3_NAK_Alborz
            PS_Traffic_2G=PS_Traffic_2G_NAK_Alborz
            RD2=RD2_NAK_Alborz
            PS_Traffic_3G=PS_Traffic_3G_NAK_Alborz
            RD3=RD3_NAK_Alborz
            PS_Traffic_4G=PS_Traffic_4G_NAK_Alborz
            RD4=RD4_NAK_Alborz
        if (t==1 ):
            Contractor="NAK-North"
            CS_Traffic=CS_Traffic_NAK_North
            CC=CC_NAK_North
            PS_Traffic=PS_Traffic_NAK_North
            RD=RD_NAK_North
            CS_Traffic_2G=CS_Traffic_2G_NAK_North
            CC2=CC2_NAK_North
            CS_Traffic_3G=CS_Traffic_3G_NAK_North
            CC3=CC3_NAK_North
            PS_Traffic_2G=PS_Traffic_2G_NAK_North
            RD2=RD2_NAK_North
            PS_Traffic_3G=PS_Traffic_3G_NAK_North
            RD3=RD3_NAK_North
            PS_Traffic_4G=PS_Traffic_4G_NAK_North
            RD4=RD4_NAK_North
        if (t==2 ):
            Contractor="NAK-Tehran"
            CS_Traffic=CS_Traffic_NAK_Tehran
            CC=CC_NAK_Tehran
            PS_Traffic=PS_Traffic_NAK_Tehran
            RD=RD_NAK_Tehran
            CS_Traffic_2G=CS_Traffic_2G_NAK_Tehran
            CC2=CC2_NAK_Tehran
            CS_Traffic_3G=CS_Traffic_3G_NAK_Tehran
            CC3=CC3_NAK_Tehran
            PS_Traffic_2G=PS_Traffic_2G_NAK_Tehran
            RD2=RD2_NAK_Tehran
            PS_Traffic_3G=PS_Traffic_3G_NAK_Tehran
            RD3=RD3_NAK_Tehran
            PS_Traffic_4G=PS_Traffic_4G_NAK_Tehran
            RD4=RD4_NAK_Tehran
        if (t==3 ):
            Contractor="NAK-Huawei"
            CS_Traffic=CS_Traffic_NAK_Huawei
            CC=CC_NAK_Huawei
            PS_Traffic=PS_Traffic_NAK_Huawei
            RD=RD_NAK_Huawei
            CS_Traffic_2G=CS_Traffic_2G_NAK_Huawei
            CC2=CC2_NAK_Huawei
            CS_Traffic_3G=CS_Traffic_3G_NAK_Huawei
            CC3=CC3_NAK_Huawei
            PS_Traffic_2G=PS_Traffic_2G_NAK_Huawei
            RD2=RD2_NAK_Huawei
            PS_Traffic_3G=PS_Traffic_3G_NAK_Huawei
            RD3=RD3_NAK_Huawei
            PS_Traffic_4G=PS_Traffic_4G_NAK_Huawei
            RD4=RD4_NAK_Huawei
        if (t==4 ):
            Contractor="NAK-Nokia"
            CS_Traffic=CS_Traffic_NAK_Nokia
            CC=CC_NAK_Nokia
            PS_Traffic=PS_Traffic_NAK_Nokia
            RD=RD_NAK_Nokia
            CS_Traffic_2G=CS_Traffic_2G_NAK_Nokia
            CC2=CC2_NAK_Nokia
            CS_Traffic_3G=CS_Traffic_3G_NAK_Nokia
            CC3=CC3_NAK_Nokia
            PS_Traffic_2G=PS_Traffic_2G_NAK_Nokia
            RD2=RD2_NAK_Nokia
            PS_Traffic_3G=PS_Traffic_3G_NAK_Nokia
            RD3=RD3_NAK_Nokia
            PS_Traffic_4G=PS_Traffic_4G_NAK_Nokia
            RD4=RD4_NAK_Nokia
        if (t==5 ):
            Contractor="BR-TEL"
            CS_Traffic=CS_Traffic_BR_TEL
            CC=CC_BR_TEL
            PS_Traffic=PS_Traffic_BR_TEL
            RD=RD_BR_TEL
            CS_Traffic_2G=CS_Traffic_2G_BR_TEL
            CC2=CC2_BR_TEL
            CS_Traffic_3G=CS_Traffic_3G_BR_TEL
            CC3=CC3_BR_TEL
            PS_Traffic_2G=PS_Traffic_2G_BR_TEL
            RD2=RD2_BR_TEL
            PS_Traffic_3G=PS_Traffic_3G_BR_TEL
            RD3=RD3_BR_TEL
            PS_Traffic_4G=PS_Traffic_4G_BR_TEL
            RD4=RD4_BR_TEL
        if (t==6 ):
            Contractor="Farafan"
            CS_Traffic=CS_Traffic_Farafan
            CC=CC_Farafan
            PS_Traffic=PS_Traffic_Farafan
            RD=RD_Farafan
            CS_Traffic_2G=CS_Traffic_2G_Farafan
            CC2=CC2_Farafan
            CS_Traffic_3G=CS_Traffic_3G_Farafan
            CC3=CC3_Farafan
            PS_Traffic_2G=PS_Traffic_2G_Farafan
            RD2=RD2_Farafan
            PS_Traffic_3G=PS_Traffic_3G_Farafan
            RD3=RD3_Farafan
            PS_Traffic_4G=PS_Traffic_4G_Farafan
            RD4=RD4_Farafan
        if (t==7 ):
            Contractor="Huawei"
            CS_Traffic=CS_Traffic_Huawei
            CC=CC_Huawei
            PS_Traffic=PS_Traffic_Huawei
            RD=RD_Huawei
            CS_Traffic_2G=CS_Traffic_2G_Huawei
            CC2=CC2_Huawei
            CS_Traffic_3G=CS_Traffic_3G_Huawei
            CC3=CC3_Huawei
            PS_Traffic_2G=PS_Traffic_2G_Huawei
            RD2=RD2_Huawei
            PS_Traffic_3G=PS_Traffic_3G_Huawei
            RD3=RD3_Huawei
            PS_Traffic_4G=PS_Traffic_4G_Huawei
            RD4=RD4_Huawei

        x = np.arange(len(CS_Traffic))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,CS_Traffic,color = "bisque")
        ax2.plot(Week_Vec,CC,color = "darkorange")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("CC-"+Contractor, fontdict = font1)
        ax1.legend(['CC(%)    Tatal Traffic (MErlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('CC_'+Contractor+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_"+Contractor+".png")
        y=30
        x=5
        h=1000
        w=550
        CC_Cropped = image[x:w, y:h]
        cv2.imwrite("CC_"+Contractor+".png", CC_Cropped)




        x = np.arange(len(PS_Traffic))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic,color = "limegreen")
        ax2.plot(Week_Vec,RD,color = "darkgreen")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD-"+Contractor, fontdict = font1)
        ax1.legend(['RD(%)    Tatal Payload (PB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD_'+Contractor+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD_"+Contractor+".png")
        y=30
        x=5
        h=1000
        w=550
        RD_Cropped = image[x:w, y:h]
        cv2.imwrite("RD_"+Contractor+".png", RD_Cropped)



        x = np.arange(len(CS_Traffic_2G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,CS_Traffic_2G,color = "lightsteelblue")
        ax2.plot(Week_Vec,CC2,color = "royalblue")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("CC2-"+Contractor, fontdict = font1)
        ax1.legend(['CC2(%)    2G Traffic (KErlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('CC2_'+Contractor+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC2_"+Contractor+".png")
        y=30
        x=5
        h=1000
        w=550
        CC2_Cropped = image[x:w, y:h]
        cv2.imwrite("CC2_"+Contractor+".png", CC2_Cropped)




        x = np.arange(len(CS_Traffic_3G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,CS_Traffic_3G,color = "khaki")
        ax2.plot(Week_Vec,CC3,color = "darkkhaki")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("CC3-"+Contractor, fontdict = font1)
        ax1.legend(['CC3(%)    3G Traffic (KErlang)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('CC3_'+Contractor+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC3_"+Contractor+".png")
        y=30
        x=5
        h=1000
        w=550
        CC3_Cropped = image[x:w, y:h]
        cv2.imwrite("CC3_"+Contractor+".png", CC3_Cropped)


        x = np.arange(len(PS_Traffic_2G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic_2G,color = "lightgrey")
        ax2.plot(Week_Vec,RD2,color = "dimgrey")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD2-"+Contractor, fontdict = font1)
        ax1.legend(['RD2(%)    2G Payload (TB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD2_'+Contractor+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD2_"+Contractor+".png")
        y=30
        x=5
        h=1000
        w=550
        RD2_Cropped = image[x:w, y:h]
        cv2.imwrite("RD2_"+Contractor+".png", RD2_Cropped)





        x = np.arange(len(PS_Traffic_3G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic_3G,color = "thistle")
        ax2.plot(Week_Vec,RD3,color = "purple")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD3-"+Contractor, fontdict = font1)
        ax1.legend(['RD3(%)    3G Payload (TB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD3_'+Contractor+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD3_"+Contractor+".png")
        y=30
        x=5
        h=1000
        w=550
        RD3_Cropped = image[x:w, y:h]
        cv2.imwrite("RD3_"+Contractor+".png", RD3_Cropped)




        x = np.arange(len(PS_Traffic_4G))
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(13),cm_to_inch(6.5)))
        ax2 = ax1.twinx()
        ax1.yaxis.tick_right()
        ax2.yaxis.tick_left()
        ax1.bar(Week_Vec,PS_Traffic_4G,color = "salmon")
        ax2.plot(Week_Vec,RD4,color = "darkred")
        ax1.set_xticks(x, Week_Vec,fontsize=5, rotation='vertical')
        font1 = {'family':'serif','color':'black','size':8}
        plt.title("RD4-"+Contractor, fontdict = font1)
        ax1.legend(['RD4(%)    4G Payload (TB)'],fontsize=7,loc='upper center',fancybox=True, shadow=True)
        ax1.yaxis.set_tick_params(labelsize=7)
        ax2.yaxis.set_tick_params(labelsize=7)
        plt.savefig('RD4_'+Contractor+'.png')

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD4_"+Contractor+".png")
        y=30
        x=5
        h=1000
        w=550
        RD4_Cropped = image[x:w, y:h]
        cv2.imwrite("RD4_"+Contractor+".png", RD4_Cropped)





        slide = prs.slides.add_slide(blank_slide_layout)

        pic_left_1  = int(prs.slide_width *0)
        pic_top_1   = int(prs.slide_width *0.02)
        pic_width_1 = int(prs.slide_width *0.5)

        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        pic_left_1  = int(prs.slide_width *0.5)
        pic_top_1   = int(prs.slide_width *0.02)
        pic_width_1 = int(prs.slide_width *0.5)

        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)


        pic_left_1  = int(prs.slide_width *0)
        pic_top_1   = int(prs.slide_width *0.29)
        pic_width_1 = int(prs.slide_width *0.5)

        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC2_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        pic_left_1  = int(prs.slide_width *0.5)
        pic_top_1   = int(prs.slide_width *0.29)
        pic_width_1 = int(prs.slide_width *0.5)

        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC3_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)



        pic_left_1  = int(prs.slide_width *0)
        pic_top_1   = int(prs.slide_width *0.56)
        pic_width_1 = int(prs.slide_width *0.33)

        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD2_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        pic_left_1  = int(prs.slide_width *0.33)
        pic_top_1   = int(prs.slide_width *0.56)
        pic_width_1 = int(prs.slide_width *0.33)

        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD3_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)

        pic_left_1  = int(prs.slide_width *0.67)
        pic_top_1   = int(prs.slide_width *0.56)
        pic_width_1 = int(prs.slide_width *0.33)

        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD4_"+Contractor+".png", pic_left_1, pic_top_1, pic_width_1)



        prs.save('test.pptx')


# *********************************************************************************************************
# ((((((((((((((((((((((((((((   CC and RD per Contractor and Country    ))))))))))))))))))))))))))))))))))
# *********************************************************************************************************


conn_performanceDB.execute("select Date, Contractor, sum([CC (%)]*[Total Voice Traffic (Erlang)])/sum([Total Voice Traffic (Erlang)]) as 'CC (%)'"+
                           "from Province_KPI_Score_Band_CS_Daily group by Date,  Contractor order by Date, Contractor")
CC_Hourly_Table=conn_performanceDB.fetchall()

conn_performanceDB.execute("select Day, Contractor, sum([RD (%)]*[Total Payload (GB)])/sum([Total Payload (GB)]) as 'RD (%)'"+
                           "from Province_KPI_Score_Band_PS_Daily group by Day,  Contractor order by Day, Contractor")
RD_Hourly_Table=conn_performanceDB.fetchall()




CC_Hourly_NAK_Alborz=[]
CCTime_Hourly_NAK_Alborz=[]
RD_Hourly_NAK_Alborz=[]
RDTime_Hourly_NAK_Alborz=[]

CC_Hourly_NAK_North=[]
CCTime_Hourly_NAK_North=[]
RD_Hourly_NAK_North=[]
RDTime_Hourly_NAK_North=[]

CC_Hourly_NAK_Tehran=[]
CCTime_Hourly_NAK_Tehran=[]
RD_Hourly_NAK_Tehran=[]
RDTime_Hourly_NAK_Tehran=[]

CC_Hourly_NAK_Huawei=[]
CCTime_Hourly_NAK_Huawei=[]
RD_Hourly_NAK_Huawei=[]
RDTime_Hourly_NAK_Huawei=[]

CC_Hourly_NAK_Nokia=[]
CCTime_Hourly_NAK_Nokia=[]
RD_Hourly_NAK_Nokia=[]
RDTime_Hourly_NAK_Nokia=[]

CC_Hourly_BR_TEL=[]
CCTime_Hourly_BR_TEL=[]
RD_Hourly_BR_TEL=[]
RDTime_Hourly_BR_TEL=[]

CC_Hourly_Farafan=[]
CCTime_Hourly_Farafan=[]
RD_Hourly_Farafan=[]
RDTime_Hourly_Farafan=[]

CC_Hourly_Huawei=[]
CCTime_Hourly_Huawei=[]
RD_Hourly_Huawei=[]
RDTime_Hourly_Huawei=[]


CC_Hourly_Iran=[]
CCTime_Hourly_Iran=[]
RD_Hourly_Iran=[]
RDTime_Hourly_Iran=[]


for i in range(len(CC_Hourly_Table)):

            Row_Data=str(CC_Hourly_Table[i])
            Row_Data=Row_Data.split(", ")
    
            Year1=Row_Data[0]
            Date=Year1[19:23]+"/"+Row_Data[1]+"/"+Row_Data[2]
            Contractor=Row_Data[5]
            Contractor=Contractor[1:len(Contractor)-1]
            CC=Row_Data[6]
            if CC=='None)':
                continue
            CC=round(float(CC[0:len(CC)-1]),3)
            if Contractor=="NAK-Alborz":
                CC_Hourly_NAK_Alborz.append(CC)
                CCTime_Hourly_NAK_Alborz.append(Date)
            if Contractor=="NAK-North":
                CC_Hourly_NAK_North.append(CC)
                CCTime_Hourly_NAK_North.append(Date)
            if Contractor=="NAK-Tehran":
                CC_Hourly_NAK_Tehran.append(CC)
                CCTime_Hourly_NAK_Tehran.append(Date)
            if Contractor=="NAK-Huawei":
                CC_Hourly_NAK_Huawei.append(CC)
                CCTime_Hourly_NAK_Huawei.append(Date)
            if Contractor=="NAK-Nokia":
                CC_Hourly_NAK_Nokia.append(CC)
                CCTime_Hourly_NAK_Nokia.append(Date)
            if Contractor=="BR-TEL":
                CC_Hourly_BR_TEL.append(CC)
                CCTime_Hourly_BR_TEL.append(Date)
            if Contractor=="Farafan":
                CC_Hourly_Farafan.append(CC)
                CCTime_Hourly_Farafan.append(Date)
            if Contractor=="Huawei":
                CC_Hourly_Huawei.append(CC)
                CCTime_Hourly_Huawei.append(Date)
            if Contractor=="IRAN":
                CC_Hourly_Iran.append(CC)
                CCTime_Hourly_Iran.append(Date)

for t in range(8):
        if (t==0 ):
            Contractor="NAK-Alborz"
            CC_Hourly=CC_Hourly_NAK_Alborz
            CCTime_Hourly=CCTime_Hourly_NAK_Alborz
        if (t==1 ):
            Contractor="NAK-North"
            CC_Hourly=CC_Hourly_NAK_North
            CCTime_Hourly=CCTime_Hourly_NAK_North
        if (t==2 ):
            Contractor="NAK-Tehran"
            CC_Hourly=CC_Hourly_NAK_Tehran
            CCTime_Hourly=CCTime_Hourly_NAK_Tehran
        if (t==3 ):
            Contractor="NAK-Huawei"
            CC_Hourly=CC_Hourly_NAK_Huawei
            CCTime_Hourly=CCTime_Hourly_NAK_Huawei
        if (t==4 ):
            Contractor="NAK-Nokia"
            CC_Hourly=CC_Hourly_NAK_Nokia
            CCTime_Hourly=CCTime_Hourly_NAK_Nokia
        if (t==5 ):
            Contractor="BR-TEL"
            CC_Hourly=CC_Hourly_BR_TEL
            CCTime_Hourly=CCTime_Hourly_BR_TEL
        if (t==6 ):
            Contractor="Farafan"
            CC_Hourly=CC_Hourly_Farafan
            CCTime_Hourly=CCTime_Hourly_Farafan
        if (t==7 ):
            Contractor="Huawei"
            CC_Hourly=CC_Hourly_Huawei
            CCTime_Hourly=CCTime_Hourly_Huawei

        downsample_Rate=round(len(CCTime_Hourly)/50)
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(32),cm_to_inch(11)))
        x_Downsample=downsample(CCTime_Hourly,downsample_Rate)
        X_Vec=[]
        x_index=0
        while len(X_Vec)!=len(x_Downsample):
             X_Vec.append(x_index)
             x_index=x_index+downsample_Rate
        ax1.plot(CCTime_Hourly, CC_Hourly, label=Contractor)
        ax1.plot(CCTime_Hourly_Iran, CC_Hourly_Iran,color = "darkorange", label='Country')
        ax1.set_xticks(X_Vec, x_Downsample,fontsize=5, rotation='vertical')
        leg = ax1.legend();
        font1 = {'family':'serif','color':'black','size':12}
        plt.title(Contractor+" & Country", fontdict = font1)
        plt.ylabel('CC (%)')
        grid(True)
        plt.savefig("CC_"+Contractor+"_Country.png")

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_"+Contractor+"_Country.png")
        y=100
        x=20
        h=1150
        w=1050
        CC_Cropped = image[x:w, y:h]
        cv2.imwrite("CC_"+Contractor+"_Country.png", CC_Cropped)


for i in range(len(RD_Hourly_Table)):

            Row_Data=str(RD_Hourly_Table[i])
            Row_Data=Row_Data.split(", ")
    
            Year1=Row_Data[0]
            Date=Year1[19:23]+"/"+Row_Data[1]+"/"+Row_Data[2]
            Contractor=Row_Data[5]
            Contractor=Contractor[1:len(Contractor)-1]
            RD=Row_Data[6]
            if RD=='None)':
                continue
            RD=round(float(RD[0:len(RD)-1]),3)
            if Contractor=="NAK-Alborz":
                RD_Hourly_NAK_Alborz.append(RD)
                RDTime_Hourly_NAK_Alborz.append(Date)
            if Contractor=="NAK-North":
                RD_Hourly_NAK_North.append(RD)
                RDTime_Hourly_NAK_North.append(Date)
            if Contractor=="NAK-Tehran":
                RD_Hourly_NAK_Tehran.append(RD)
                RDTime_Hourly_NAK_Tehran.append(Date)
            if Contractor=="NAK-Huawei":
                RD_Hourly_NAK_Huawei.append(RD)
                RDTime_Hourly_NAK_Huawei.append(Date)
            if Contractor=="NAK-Nokia":
                RD_Hourly_NAK_Nokia.append(RD)
                RDTime_Hourly_NAK_Nokia.append(Date)
            if Contractor=="BR-TEL":
                RD_Hourly_BR_TEL.append(RD)
                RDTime_Hourly_BR_TEL.append(Date)
            if Contractor=="Farafan":
                RD_Hourly_Farafan.append(RD)
                RDTime_Hourly_Farafan.append(Date)
            if Contractor=="Huawei":
                RD_Hourly_Huawei.append(RD)
                RDTime_Hourly_Huawei.append(Date)
            if Contractor=="IRAN":
                RD_Hourly_Iran.append(RD)
                RDTime_Hourly_Iran.append(Date)

for t in range(8):
        if (t==0 ):
            Contractor="NAK-Alborz"
            RD_Hourly=RD_Hourly_NAK_Alborz
            RDTime_Hourly=RDTime_Hourly_NAK_Alborz
        if (t==1 ):
            Contractor="NAK-North"
            RD_Hourly=RD_Hourly_NAK_North
            RDTime_Hourly=RDTime_Hourly_NAK_North
        if (t==2 ):
            Contractor="NAK-Tehran"
            RD_Hourly=RD_Hourly_NAK_Tehran
            RDTime_Hourly=RDTime_Hourly_NAK_Tehran
        if (t==3 ):
            Contractor="NAK-Huawei"
            RD_Hourly=RD_Hourly_NAK_Huawei
            RDTime_Hourly=RDTime_Hourly_NAK_Huawei
        if (t==4 ):
            Contractor="NAK-Nokia"
            RD_Hourly=RD_Hourly_NAK_Nokia
            RDTime_Hourly=RDTime_Hourly_NAK_Nokia
        if (t==5 ):
            Contractor="BR-TEL"
            RD_Hourly=RD_Hourly_BR_TEL
            RDTime_Hourly=RDTime_Hourly_BR_TEL
        if (t==6 ):
            Contractor="Farafan"
            RD_Hourly=RD_Hourly_Farafan
            RDTime_Hourly=RDTime_Hourly_Farafan
        if (t==7 ):
            Contractor="Huawei"
            RD_Hourly=RD_Hourly_Huawei
            RDTime_Hourly=RDTime_Hourly_Huawei

        downsample_Rate=round(len(RDTime_Hourly)/50)
        fig, ax1 = plt.subplots(figsize=(cm_to_inch(32),cm_to_inch(11)))
        x_Downsample=downsample(RDTime_Hourly,downsample_Rate)
        X_Vec=[]
        x_index=0
        while len(X_Vec)!=len(x_Downsample):
             X_Vec.append(x_index)
             x_index=x_index+downsample_Rate
        ax1.plot(RDTime_Hourly, RD_Hourly, label=Contractor)
        ax1.plot(RDTime_Hourly_Iran, RD_Hourly_Iran,color = "darkorange", label='Country')
        ax1.set_xticks(X_Vec, x_Downsample,fontsize=5, rotation='vertical')
        leg = ax1.legend();
        font1 = {'family':'serif','color':'black','size':12}
        plt.title(Contractor+" & Country", fontdict = font1)
        plt.ylabel('RD (%)')
        grid(True)
        plt.savefig("RD_"+Contractor+"_Country.png")

        image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD_"+Contractor+"_Country.png")
        y=100
        x=20
        h=1150
        w=1050
        RD_Cropped = image[x:w, y:h]
        cv2.imwrite("RD_"+Contractor+"_Country.png", RD_Cropped)

        
        slide = prs.slides.add_slide(blank_slide_layout)

        pic_left_1  = int(prs.slide_width *0.047)
        pic_top_1   = int(prs.slide_width *0.01)
        pic_width_1 = int(prs.slide_width *0.9)


        pic_left_3  = int(prs.slide_width *0.047)
        pic_top_3   = int(prs.slide_width *0.38)
        pic_width_3 = int(prs.slide_width *0.9)



        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC_"+Contractor+"_Country.png", pic_left_1, pic_top_1, pic_width_1)
        pic = slide.shapes.add_picture("D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\RD_"+Contractor+"_Country.png", pic_left_3, pic_top_3, pic_width_3)


        prs.save('test.pptx')



# ******************************************************************************************
# (((((((((((((((((((((((((((((((((((((   CS and PS    )))))))))))))))))))))))))))))))))))))
# ******************************************************************************************

conn_performanceDB.execute("select Wk,Contractor, avg([CSSR_2G]) as 'CSSR_2G', avg([CSSR_3G]) as 'CSSR_3G', avg([CDR_2G]) as 'CDR_2G', avg([CDR_3G]) as 'CDR_3G' from ("+
                           "select Wk, Contractor,PIndex, SUM([2G TCH Traffic]*[CSSR_MCI])/sum([2G TCH Traffic]) as 'CSSR_2G',"+
                              "SUM([3G_CS_Traffic]*[CS_CSSR])/sum([3G_CS_Traffic]) as 'CSSR_3G',"+
                              "SUM([2G TCH Traffic]*[CDR])/sum([2G TCH Traffic]) as 'CDR_2G',"+
                              "SUM([3G_CS_Traffic]*[3G_CS_Drop])/sum([3G_CS_Traffic]) as 'CDR_3G' "+
							  "from  Province_KPI_Score_Band_CS_Daily  group by Wk, Contractor, PIndex ) tble  group by Wk, Contractor  order by Wk")
CC_Table=conn_performanceDB.fetchall()


conn_performanceDB.execute("select Wk,Contractor, avg([PSSR_2G]) as 'PSSR_2G', avg([PSSR_3G]) as 'PSSR_3G', avg([PSSR_4G]) as 'PSSR_4G', avg([PDR_2G]) as 'PDR_2G', avg([PDR_3G]) as 'PDR_3G', avg([PDR_4G]) as 'PDR_4G' from ("+
                             "select Wk, Contractor,[Province Index], SUM([2G PS Traffic (GB)]*[TBF_Establishment_SR])/sum([2G PS Traffic (GB)]) as 'PSSR_2G',"+
							  "SUM([3G Payload (GB)]*[PS_CSSR])/sum([3G Payload (GB)]) as 'PSSR_3G',"+
                              "SUM([4G Payload (GB)]*[Initital E-RAB SR])/sum([4G Payload (GB)]) as 'PSSR_4G',"+
                              "SUM([2G PS Traffic (GB)]*[TBF_Drop])/sum([2G PS Traffic (GB)]) as 'PDR_2G',"+
                              "SUM([3G Payload (GB)]*[PS_Call_Drop])/sum([3G Payload (GB)]) as 'PDR_3G',"+
                              "SUM([4G Payload (GB)]*[ERAB_Drop_Rate])/sum([4G Payload (GB)]) as 'PDR_4G' "+
							  "from  Province_KPI_Score_Band_PS_Daily group by Wk, Contractor, [Province Index] ) tble group by Wk, Contractor  order by Wk")
RD_Table=conn_performanceDB.fetchall()


# Change Code from this point

CC2_NAK_Alborz=[]
CC2_NAK_Tehran=[]
CC2_NAK_North=[]
CC2_NAK_Nokia=[]
CC2_NAK_Huawei=[]
CC2_Farafan=[]
CC2_BR_TEL=[]
CC2_Huawei=[]
CC2_Iran=[]

CC3_NAK_Alborz=[]
CC3_NAK_Tehran=[]
CC3_NAK_North=[]
CC3_NAK_Nokia=[]
CC3_NAK_Huawei=[]
CC3_Farafan=[]
CC3_BR_TEL=[]
CC3_Huawei=[]
CC3_Iran=[]

CC_NAK_Alborz=[]
CC_NAK_Tehran=[]
CC_NAK_North=[]
CC_NAK_Nokia=[]
CC_NAK_Huawei=[]
CC_Farafan=[]
CC_BR_TEL=[]
CC_Huawei=[]
CC_Iran=[]



CS_Traffic_2G_NAK_Alborz=[]
CS_Traffic_2G_NAK_Tehran=[]
CS_Traffic_2G_NAK_North=[]
CS_Traffic_2G_NAK_Nokia=[]
CS_Traffic_2G_NAK_Huawei=[]
CS_Traffic_2G_Farafan=[]
CS_Traffic_2G_BR_TEL=[]
CS_Traffic_2G_Huawei=[]
CS_Traffic_2G_Iran=[]

CS_Traffic_3G_NAK_Alborz=[]
CS_Traffic_3G_NAK_Tehran=[]
CS_Traffic_3G_NAK_North=[]
CS_Traffic_3G_NAK_Nokia=[]
CS_Traffic_3G_NAK_Huawei=[]
CS_Traffic_3G_Farafan=[]
CS_Traffic_3G_BR_TEL=[]
CS_Traffic_3G_Huawei=[]
CS_Traffic_3G_Iran=[]

CS_Traffic_4G_NAK_Alborz=[]
CS_Traffic_4G_NAK_Tehran=[]
CS_Traffic_4G_NAK_North=[]
CS_Traffic_4G_NAK_Nokia=[]
CS_Traffic_4G_NAK_Huawei=[]
CS_Traffic_4G_Farafan=[]
CS_Traffic_4G_BR_TEL=[]
CS_Traffic_4G_Huawei=[]
CS_Traffic_4G_Iran=[]

Week_Vec=[]

for i in range(len(CC_Table)):
    Row_Data=str(CC_Table[i])
    Row_Data=Row_Data.split(", ")

    Week=Row_Data[0]
    Contractor=Row_Data[1]
    CC2_Str=Row_Data[2]
    CC2_Val=round(float(CC2_Str[0:len(CC2_Str)-1]),2)
    CC3_Str=Row_Data[3]
    CC3_Val=round(float(CC3_Str[0:len(CC3_Str)-1]),2)
    CC_Str=Row_Data[4]
    CC_Val=round(float(CC_Str[0:len(CC_Str)-1]),2)
    CS_2G_Str=Row_Data[5]
    CS_2G_Val=round(float(CS_2G_Str[0:len(CS_2G_Str)-1]),2)
    CS_3G_Str=Row_Data[6]
    CS_3G_Val=round(float(CS_3G_Str[0:len(CS_3G_Str)-1]),2)
    CS_4G_Str=Row_Data[7]
    CS_4G_Val=round(float(CS_4G_Str[0:len(CS_4G_Str)-1]),2)
    CS_Str=Row_Data[8]
    CS_Val=round(float(CS_Str[0:len(CS_Str)-1]),2)

    Week=Week[2:9]
    Contractor=Contractor[1:len(Contractor)-1]
    

    #if Week=='1401-33':
    #    break

    if (Contractor=='NAK-Alborz'):
        Week_Vec.append('W'+Week[5:7])
        CC2_NAK_Alborz.append(CC2_Val)
        CC3_NAK_Alborz.append(CC3_Val)
        CC_NAK_Alborz.append(CC_Val)
        CS_Traffic_2G_NAK_Alborz.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Alborz.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Alborz.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-Tehran'):
        CC2_NAK_Tehran.append(CC2_Val)
        CC3_NAK_Tehran.append(CC3_Val)
        CC_NAK_Tehran.append(CC_Val)
        CS_Traffic_2G_NAK_Tehran.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Tehran.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Tehran.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-North'):
        CC2_NAK_North.append(CC2_Val)
        CC3_NAK_North.append(CC3_Val)
        CC_NAK_North.append(CC_Val)
        CS_Traffic_2G_NAK_North.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_North.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_North.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-Nokia'):
        CC2_NAK_Nokia.append(CC2_Val)
        CC3_NAK_Nokia.append(CC3_Val)
        CC_NAK_Nokia.append(CC_Val)
        CS_Traffic_2G_NAK_Nokia.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Nokia.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Nokia.append(CS_4G_Val/1e3)
    if (Contractor=='NAK-Huawei'):
        CC2_NAK_Huawei.append(CC2_Val)
        CC3_NAK_Huawei.append(CC3_Val)
        CC_NAK_Huawei.append(CC_Val)
        CS_Traffic_2G_NAK_Huawei.append(CS_2G_Val/1e3)
        CS_Traffic_3G_NAK_Huawei.append(CS_3G_Val/1e3)
        CS_Traffic_4G_NAK_Huawei.append(CS_4G_Val/1e3)
    if (Contractor=='Farafan'):
        CC2_Farafan.append(CC2_Val)
        CC3_Farafan.append(CC3_Val)
        CC_Farafan.append(CC_Val)
        CS_Traffic_2G_Farafan.append(CS_2G_Val/1e3)
        CS_Traffic_3G_Farafan.append(CS_3G_Val/1e3)
        CS_Traffic_4G_Farafan.append(CS_4G_Val/1e3)
    if (Contractor=='BR-TEL'):
        CC2_BR_TEL.append(CC2_Val)
        CC3_BR_TEL.append(CC3_Val)
        CC_BR_TEL.append(CC_Val)
        CS_Traffic_2G_BR_TEL.append(CS_2G_Val/1e3)
        CS_Traffic_3G_BR_TEL.append(CS_3G_Val/1e3)
        CS_Traffic_4G_BR_TEL.append(CS_4G_Val/1e3)
    if (Contractor=='Huawei'):
        CC2_Huawei.append(CC2_Val)
        CC3_Huawei.append(CC3_Val)
        CC_Huawei.append(CC_Val)
        CS_Traffic_2G_Huawei.append(CS_2G_Val/1e3)
        CS_Traffic_3G_Huawei.append(CS_3G_Val/1e3)
        CS_Traffic_4G_Huawei.append(CS_4G_Val/1e3)
    if (Contractor=='IRAN'):
        CC2_Iran.append(CC2_Val)
        CC3_Iran.append(CC3_Val)
        CC_Iran.append(CC_Val)
        CS_Traffic_2G_Iran.append(CS_2G_Val/1e3)
        CS_Traffic_3G_Iran.append(CS_3G_Val/1e3)
        CS_Traffic_4G_Iran.append(CS_4G_Val/1e3)



Last_CC2=[CC2_NAK_Alborz[len(CC2_NAK_Alborz)-1], CC2_NAK_Tehran[len(CC2_NAK_Alborz)-1], CC2_NAK_North[len(CC2_NAK_Alborz)-1], CC2_NAK_Nokia[len(CC2_NAK_Alborz)-1], CC2_NAK_Huawei[len(CC2_NAK_Alborz)-1], CC2_Farafan[len(CC2_NAK_Alborz)-1], CC2_BR_TEL[len(CC2_NAK_Alborz)-1], CC2_Huawei[len(CC2_NAK_Alborz)-1] ]
Index_of_Sort_CC2=np.argsort(Last_CC2)
Data_Sorted_Array_CC2=[]
x_Labels_CC2=[];
for k in range(len(Index_of_Sort_CC2)):
    if Index_of_Sort_CC2[k]==0:
        Data_Sorted_Array_CC2.append(CC2_NAK_Alborz)
        x_Labels_CC2.append('NAK-Alborz')
    if Index_of_Sort_CC2[k]==1:
        Data_Sorted_Array_CC2.append(CC2_NAK_Tehran)
        x_Labels_CC2.append('NAK-Tehran')
    if Index_of_Sort_CC2[k]==2:
        Data_Sorted_Array_CC2.append(CC2_NAK_North)
        x_Labels_CC2.append('NAK-North')
    if Index_of_Sort_CC2[k]==3:
        Data_Sorted_Array_CC2.append(CC2_NAK_Nokia)
        x_Labels_CC2.append('NAK-Nokia')
    if Index_of_Sort_CC2[k]==4:
        Data_Sorted_Array_CC2.append(CC2_NAK_Huawei)
        x_Labels_CC2.append('NAK-Huawei')
    if Index_of_Sort_CC2[k]==5:
        Data_Sorted_Array_CC2.append(CC2_Farafan)
        x_Labels_CC2.append('Farafan')
    if Index_of_Sort_CC2[k]==6:
        Data_Sorted_Array_CC2.append(CC2_BR_TEL)
        x_Labels_CC2.append('BR-TEL')
    if Index_of_Sort_CC2[k]==7:
        Data_Sorted_Array_CC2.append(CC2_Huawei)
        x_Labels_CC2.append('Huawei')


data=np.array(Data_Sorted_Array_CC2)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_CC2):
    plt.text( i + dx[31],Last_CC2[Index_of_Sort_CC2[i]] , str(Last_CC2[Index_of_Sort_CC2[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_CC2)
font1 = {'family':'serif','color':'black','size':17}
plt.title("CC2(%)", fontdict = font1)
plt.ylim(75, 100)
grid(True)
plt.savefig('CC2.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC2.png")
y=80
x=10
h=1000
w=520
CC2_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("CC2.png", CC2_Bar_Cropped)



Last_CC3=[CC3_NAK_Alborz[len(CC3_NAK_Alborz)-1], CC3_NAK_Tehran[len(CC3_NAK_Alborz)-1], CC3_NAK_North[len(CC3_NAK_Alborz)-1], CC3_NAK_Nokia[len(CC3_NAK_Alborz)-1], CC3_NAK_Huawei[len(CC3_NAK_Alborz)-1], CC3_Farafan[len(CC3_NAK_Alborz)-1], CC3_BR_TEL[len(CC3_NAK_Alborz)-1], CC3_Huawei[len(CC3_NAK_Alborz)-1] ]
Index_of_Sort_CC3=np.argsort(Last_CC3)
Data_Sorted_Array_CC3=[]
x_Labels_CC3=[];
for k in range(len(Index_of_Sort_CC3)):
    if Index_of_Sort_CC3[k]==0:
        Data_Sorted_Array_CC3.append(CC3_NAK_Alborz)
        x_Labels_CC3.append('NAK-Alborz')
    if Index_of_Sort_CC3[k]==1:
        Data_Sorted_Array_CC3.append(CC3_NAK_Tehran)
        x_Labels_CC3.append('NAK-Tehran')
    if Index_of_Sort_CC3[k]==2:
        Data_Sorted_Array_CC3.append(CC3_NAK_North)
        x_Labels_CC3.append('NAK-North')
    if Index_of_Sort_CC3[k]==3:
        Data_Sorted_Array_CC3.append(CC3_NAK_Nokia)
        x_Labels_CC3.append('NAK-Nokia')
    if Index_of_Sort_CC3[k]==4:
        Data_Sorted_Array_CC3.append(CC3_NAK_Huawei)
        x_Labels_CC3.append('NAK-Huawei')
    if Index_of_Sort_CC3[k]==5:
        Data_Sorted_Array_CC3.append(CC3_Farafan)
        x_Labels_CC3.append('Farafan')
    if Index_of_Sort_CC3[k]==6:
        Data_Sorted_Array_CC3.append(CC3_BR_TEL)
        x_Labels_CC3.append('BR-TEL')
    if Index_of_Sort_CC3[k]==7:
        Data_Sorted_Array_CC3.append(CC3_Huawei)
        x_Labels_CC3.append('Huawei')


data=np.array(Data_Sorted_Array_CC3)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_CC3):
    plt.text( i + dx[31],Last_CC3[Index_of_Sort_CC3[i]] , str(Last_CC3[Index_of_Sort_CC3[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_CC3)
font1 = {'family':'serif','color':'black','size':17}
plt.title("CC3(%)", fontdict = font1)
plt.ylim(75, 100)
grid(True)
plt.savefig('CC3.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC3.png")
y=80
x=10
h=1000
w=520
CC3_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("CC3.png", CC3_Bar_Cropped)

Last_CC=[CC_NAK_Alborz[len(CC_NAK_Alborz)-1], CC_NAK_Tehran[len(CC_NAK_Alborz)-1], CC_NAK_North[len(CC_NAK_Alborz)-1], CC_NAK_Nokia[len(CC_NAK_Alborz)-1], CC_NAK_Huawei[len(CC_NAK_Alborz)-1], CC_Farafan[len(CC_NAK_Alborz)-1], CC_BR_TEL[len(CC_NAK_Alborz)-1], CC_Huawei[len(CC_NAK_Alborz)-1] ]
Index_of_Sort_CC=np.argsort(Last_CC)
Data_Sorted_Array_CC=[]
x_Labels_CC=[];
for k in range(len(Index_of_Sort_CC)):
    if Index_of_Sort_CC[k]==0:
        Data_Sorted_Array_CC.append(CC_NAK_Alborz)
        x_Labels_CC.append('NAK-Alborz')
    if Index_of_Sort_CC[k]==1:
        Data_Sorted_Array_CC.append(CC_NAK_Tehran)
        x_Labels_CC.append('NAK-Tehran')
    if Index_of_Sort_CC[k]==2:
        Data_Sorted_Array_CC.append(CC_NAK_North)
        x_Labels_CC.append('NAK-North')
    if Index_of_Sort_CC[k]==3:
        Data_Sorted_Array_CC.append(CC_NAK_Nokia)
        x_Labels_CC.append('NAK-Nokia')
    if Index_of_Sort_CC[k]==4:
        Data_Sorted_Array_CC.append(CC_NAK_Huawei)
        x_Labels_CC.append('NAK-Huawei')
    if Index_of_Sort_CC[k]==5:
        Data_Sorted_Array_CC.append(CC_Farafan)
        x_Labels_CC.append('Farafan')
    if Index_of_Sort_CC[k]==6:
        Data_Sorted_Array_CC.append(CC_BR_TEL)
        x_Labels_CC.append('BR-TEL')
    if Index_of_Sort_CC[k]==7:
        Data_Sorted_Array_CC.append(CC_Huawei)
        x_Labels_CC.append('Huawei')


data=np.array(Data_Sorted_Array_CC)

x = np.arange(data.shape[0])
dx = (np.arange(data.shape[1])-data.shape[1]/2.)/(data.shape[1]+2.)
d = 1./(data.shape[1]+2.)

def cm_to_inch(value):
    return value/2.54
plt.figure(figsize=(cm_to_inch(27),cm_to_inch(9)))
axes= plt.axes()


for i in range(data.shape[1]):
    plt.bar(x+dx[i],data[:,i], color = "orange",  width=d, label="label {}".format(i))


for i , v in enumerate(Last_CC):
    plt.text( i + dx[31],Last_CC[Index_of_Sort_CC[i]] , str(Last_CC[Index_of_Sort_CC[i]]), color='black', size=12, fontweight='bold')

axes.set_xticks(x, x_Labels_CC)
font1 = {'family':'serif','color':'black','size':17}
plt.title("CC(%)", fontdict = font1)
plt.ylim(75, 100)
grid(True)
plt.savefig('CC.png')

image = cv2.imread(r"D:\P1\Performane\Programmes\Python Projects\Weekly_Dashboards\Weekly_Dashboards\Weekly_Dashboards\CC.png")
y=80
x=10
h=1000
w=520
CC_Bar_Cropped = image[x:w, y:h]
cv2.imwrite("CC.png", CC_Bar_Cropped)

